"""Build the deterministic eval fixture deck (evals/fixtures/eval_deck.pptx).

The deck exercises the read/inspect/compliance surface and contains exactly
one seeded compliance violation (a non-theme font on the last slide's title).
Run: uv run python scripts/make_eval_fixtures.py
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Pt

OUT = Path(__file__).parent.parent / "evals" / "fixtures" / "eval_deck.pptx"


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    image_path = OUT.parent / "vision.png"
    Image.new("RGB", (640, 480), "#1A6BCC").save(image_path)

    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    s1.shapes.title.text = "Project Aurora Strategy"
    s1.placeholders[1].text = "Board Briefing 2026"

    s2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    s2.shapes.title.text = "Strategic Pillars"
    tf = s2.placeholders[1].text_frame
    tf.text = "Platform consolidation"
    tf.add_paragraph().text = "Data governance"
    tf.add_paragraph().text = "Talent uplift"

    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Key Risks"
    tf = s3.placeholders[1].text_frame
    tf.text = "Vendor lock-in"
    tf.add_paragraph().text = "Skills gap"
    s3.notes_slide.notes_text_frame.text = (
        "Verify the audit trail before the board meeting."
    )

    s4 = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
    s4.shapes.title.text = "Vision"
    s4.placeholders[1].insert_picture(str(image_path))

    s5 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s5.shapes.title.text = "Thank You"
    run = s5.shapes.title.text_frame.paragraphs[0].runs[0]
    run.font.name = "Comic Sans MS"  # the seeded C02 violation
    run.font.size = Pt(40)

    prs.save(str(OUT))
    print(f"saved {OUT}")


if __name__ == "__main__":
    main()

"""Compliance engine (DESIGN.md §7): rules C01–C10 over a deck, validated
against a registered template (full check) or the deck's own theme
(internal-consistency mode).

Findings are structured and severity-graded; repairs are planned per finding
and only zero-risk fixes are applied under the conservative strategy.
Text-fit checks are arithmetic estimates (confidence: estimated) until a
render confirms them.
"""

from __future__ import annotations

import math
from typing import Any

from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from .reader import emu_to_inches
from .templates import extract_theme_from_master

# Estimation constants: average glyph width ~0.55 em, line height ~1.25 em.
_CHAR_EM = 0.55
_LINE_EM = 1.25
_DEFAULT_FONT_PT = 18.0
_OVERFLOW_TOLERANCE = 1.10  # flag only when estimated need exceeds box by 10 %
_COLOR_SNAP_DISTANCE = 60.0  # max RGB distance for an auto-snap repair


def _finding(
    rule: str,
    severity: str,
    message: str,
    slide_index: int | None = None,
    shape_id: int | None = None,
    auto_fixable: bool = False,
    confidence: str = "exact",
) -> dict[str, Any]:
    return {
        "rule": rule,
        "severity": severity,
        "slide_index": slide_index,
        "shape_id": shape_id,
        "message": message,
        "auto_fixable": auto_fixable,
        "confidence": confidence,
    }


# -- helpers ----------------------------------------------------------------------


def _theme_for(prs: PresentationType, template_entry: dict | None) -> dict:
    if template_entry is not None:
        return template_entry["theme"]
    return extract_theme_from_master(prs.slide_masters[0])


def _theme_font_ok(name: str | None, fonts: dict) -> bool:
    if name is None or name.startswith("+"):  # inherited or theme reference
        return True
    return name in {fonts.get("major"), fonts.get("minor")}


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def _color_distance(a: str, b: str) -> float:
    ra, ga, ba = _hex_to_rgb(a)
    rb, gb, bb = _hex_to_rgb(b)
    return math.sqrt((ra - rb) ** 2 + (ga - gb) ** 2 + (ba - bb) ** 2)


def _nearest_theme_color(value: str, palette: dict[str, str]) -> tuple[str, str, float]:
    best_name, best_hex, best_dist = "", "", float("inf")
    for name, hex_value in palette.items():
        dist = _color_distance(value, hex_value)
        if dist < best_dist:
            best_name, best_hex, best_dist = name, hex_value, dist
    return best_name, best_hex, best_dist


def _shape_geometry(shape, slide: Slide) -> tuple[float, float, float, float] | None:
    """(left, top, width, height) in inches, resolving layout inheritance."""
    if shape.left is not None:
        return (
            emu_to_inches(shape.left),
            emu_to_inches(shape.top),
            emu_to_inches(shape.width) or 0.0,
            emu_to_inches(shape.height) or 0.0,
        )
    if shape.is_placeholder:
        ph_idx = shape.placeholder_format.idx
        for scope in (slide.slide_layout, slide.slide_layout.slide_master):
            for ph in scope.placeholders:
                if ph.placeholder_format.idx == ph_idx and ph.left is not None:
                    return (
                        emu_to_inches(ph.left),
                        emu_to_inches(ph.top),
                        emu_to_inches(ph.width) or 0.0,
                        emu_to_inches(ph.height) or 0.0,
                    )
    return None


def _overlap_fraction(a, b) -> float:
    """Fraction of box b covered by box a (boxes: left, top, width, height)."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    area_b = bw * bh
    return (ix * iy) / area_b if area_b else 0.0


def _iter_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        yield from paragraph.runs


# -- rules ---------------------------------------------------------------------


def _check_layout_provenance(prs, template_entry, findings) -> None:
    """C01: every slide's layout exists in the reference template (by name)."""
    template_layouts = {lo["name"].lower() for lo in template_entry["layouts"]}
    for i, slide in enumerate(prs.slides):
        name = slide.slide_layout.name
        if name.lower() not in template_layouts:
            findings.append(
                _finding(
                    "C01", "error",
                    f"Slide uses layout '{name}' which is not part of the reference "
                    "template. Re-create it from a template layout or use "
                    "ppt_apply_template once available.",
                    slide_index=i + 1,
                )
            )


def _check_fonts(prs, theme, findings) -> None:
    """C02: explicit run fonts must resolve to the theme font scheme."""
    fonts = theme["font_scheme"]
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            bad = {
                run.font.name
                for run in _iter_runs(shape)
                if not _theme_font_ok(run.font.name, fonts)
            }
            for name in sorted(bad):
                findings.append(
                    _finding(
                        "C02", "error",
                        f"Font '{name}' is not in the theme font scheme "
                        f"(major='{fonts.get('major')}', minor='{fonts.get('minor')}'). "
                        "Conservative repair re-links the runs to the theme font.",
                        slide_index=i + 1, shape_id=shape.shape_id, auto_fixable=True,
                    )
                )


def _check_colors(prs, theme, findings) -> None:
    """C03: explicit RGB font colors should come from the theme palette."""
    palette = theme["color_scheme"]
    if not palette:
        return
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            seen: set[str] = set()
            for run in _iter_runs(shape):
                color = run.font.color
                try:
                    if color.type != MSO_COLOR_TYPE.RGB:
                        continue
                    value = f"#{color.rgb}"
                except (AttributeError, ValueError):
                    continue
                if value.upper() in {v.upper() for v in palette.values()} or value in seen:
                    continue
                seen.add(value)
                name, hex_value, dist = _nearest_theme_color(value, palette)
                snappable = dist <= _COLOR_SNAP_DISTANCE
                findings.append(
                    _finding(
                        "C03", "warning",
                        f"Explicit color {value} is not a theme color. Nearest: "
                        f"{name} ({hex_value}, distance {dist:.0f})"
                        + (" — close enough for auto-snap." if snappable else "."),
                        slide_index=i + 1, shape_id=shape.shape_id,
                        auto_fixable=snappable,
                    )
                )


def _check_footers(prs, findings) -> None:
    """C04: layout defines footer-family placeholders the slide lacks."""
    from .writer import placeholder_role

    for i, slide in enumerate(prs.slides):
        layout_roles = {placeholder_role(ph) for ph in slide.slide_layout.placeholders}
        slide_roles = {placeholder_role(ph) for ph in slide.placeholders}
        missing = (layout_roles & {"footer", "slide_number", "date"}) - slide_roles
        if missing:
            findings.append(
                _finding(
                    "C04", "info",
                    f"Layout '{slide.slide_layout.name}' defines {sorted(missing)} "
                    "placeholder(s) this slide does not carry. PowerPoint's "
                    "'Insert > Header & Footer' normally populates these.",
                    slide_index=i + 1,
                )
            )


def _check_master_elements_covered(prs, findings) -> None:
    """C05: recurring master shapes (logos etc.) must not be covered."""
    for i, slide in enumerate(prs.slides):
        master = slide.slide_layout.slide_master
        logos = [
            (shape, geometry)
            for shape in master.shapes
            if not shape.is_placeholder
            and (geometry := _shape_geometry(shape, slide)) is not None
        ]
        for shape in slide.shapes:
            geometry = _shape_geometry(shape, slide)
            if geometry is None:
                continue
            for logo, logo_geometry in logos:
                if _overlap_fraction(geometry, logo_geometry) >= 0.8:
                    findings.append(
                        _finding(
                            "C05", "error",
                            f"Shape '{shape.name}' covers the master element "
                            f"'{logo.name}' (logo/recurring brand element).",
                            slide_index=i + 1, shape_id=shape.shape_id,
                        )
                    )


def _check_placeholder_bypass(prs, findings) -> None:
    """C06: free text shapes overlapping an *empty* placeholder."""
    for i, slide in enumerate(prs.slides):
        empty_placeholders = [
            (ph, geometry)
            for ph in slide.placeholders
            if getattr(ph, "has_text_frame", False)
            and not ph.text_frame.text.strip()
            and (geometry := _shape_geometry(ph, slide)) is not None
        ]
        if not empty_placeholders:
            continue
        for shape in slide.shapes:
            if shape.is_placeholder or not getattr(shape, "has_text_frame", False):
                continue
            if not shape.text_frame.text.strip():
                continue
            geometry = _shape_geometry(shape, slide)
            if geometry is None:
                continue
            for ph, ph_geometry in empty_placeholders:
                if _overlap_fraction(geometry, ph_geometry) >= 0.5:
                    findings.append(
                        _finding(
                            "C06", "warning",
                            f"Free textbox '{shape.name}' sits on top of the empty "
                            f"placeholder idx:{ph.placeholder_format.idx} — the "
                            "content belongs in the placeholder "
                            "(ppt_set_placeholder_content).",
                            slide_index=i + 1, shape_id=shape.shape_id,
                        )
                    )


def _estimate_text_height_in(shape, width_in: float) -> float:
    total_lines = 0
    max_pt = 0.0
    for paragraph in shape.text_frame.paragraphs:
        sizes = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
        pt = max(sizes) if sizes else _DEFAULT_FONT_PT
        max_pt = max(max_pt, pt)
        chars_per_line = max(1.0, width_in / (_CHAR_EM * pt / 72))
        text = paragraph.text
        total_lines += max(1, math.ceil(len(text) / chars_per_line)) if text else 1
    return total_lines * _LINE_EM * (max_pt or _DEFAULT_FONT_PT) / 72


def _check_overflow(prs, findings) -> None:
    """C07: estimated text extent exceeds the shape's box."""
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if not shape.text_frame.text.strip():
                continue
            geometry = _shape_geometry(shape, slide)
            if geometry is None or geometry[2] <= 0 or geometry[3] <= 0:
                continue
            needed = _estimate_text_height_in(shape, geometry[2])
            if needed > geometry[3] * _OVERFLOW_TOLERANCE:
                findings.append(
                    _finding(
                        "C07", "error",
                        f"Text in '{shape.name}' needs ~{needed:.1f} in but the box is "
                        f"{geometry[3]:.1f} in tall — likely overflow. Shorten the "
                        "content or pick a roomier layout.",
                        slide_index=i + 1, shape_id=shape.shape_id,
                        confidence="estimated",
                    )
                )


def _check_density(prs, template_entry, findings) -> None:
    """C08: bullet count vs the template layout's estimated capacity."""
    capacity_by_name = {
        lo["name"].lower(): lo["capacity"] for lo in template_entry["layouts"]
    }
    for i, slide in enumerate(prs.slides):
        capacity = capacity_by_name.get(slide.slide_layout.name.lower())
        if not capacity or not capacity.get("body_bullets"):
            continue
        from .writer import placeholder_role

        for ph in slide.placeholders:
            if placeholder_role(ph) != "body" or not getattr(ph, "has_text_frame", False):
                continue
            bullets = len([p for p in ph.text_frame.paragraphs if p.text.strip()])
            if bullets > capacity["body_bullets"]:
                findings.append(
                    _finding(
                        "C08", "warning",
                        f"{bullets} bullets exceed the layout's estimated capacity of "
                        f"{capacity['body_bullets']} — consider splitting the slide.",
                        slide_index=i + 1, shape_id=ph.shape_id,
                        confidence="estimated",
                    )
                )


def _check_offgrid(prs, findings) -> None:
    """C09: freeform shapes not aligned to any layout placeholder edge."""
    tolerance = 0.1
    for i, slide in enumerate(prs.slides):
        anchors: list[float] = []
        for ph in slide.slide_layout.placeholders:
            if ph.left is None:
                continue
            anchors.extend(
                [emu_to_inches(ph.left), emu_to_inches(ph.top),
                 emu_to_inches(ph.left) + (emu_to_inches(ph.width) or 0)]
            )
        if not anchors:
            continue
        for shape in slide.shapes:
            if shape.is_placeholder or shape.left is None:
                continue
            left = emu_to_inches(shape.left)
            if all(abs(left - anchor) > tolerance for anchor in anchors):
                findings.append(
                    _finding(
                        "C09", "info",
                        f"Manually positioned shape '{shape.name}' (left={left} in) "
                        "does not align with the layout grid.",
                        slide_index=i + 1, shape_id=shape.shape_id,
                        confidence="estimated",
                    )
                )


def _check_slide_size(prs, template_entry, findings) -> None:
    """C10: slide size must match the template."""
    expected = template_entry["slide_size"]
    actual_w = emu_to_inches(int(prs.slide_width))
    actual_h = emu_to_inches(int(prs.slide_height))
    if (actual_w, actual_h) != (expected["width_in"], expected["height_in"]):
        findings.append(
            _finding(
                "C10", "error",
                f"Slide size {actual_w} x {actual_h} in differs from the template's "
                f"{expected['width_in']} x {expected['height_in']} in.",
            )
        )


def validate(prs: PresentationType, template_entry: dict | None) -> list[dict[str, Any]]:
    """Run all applicable rules. Without a template entry, template-relative
    rules (C01, C08, C10) are skipped and the deck's own theme is the
    reference."""
    findings: list[dict[str, Any]] = []
    theme = _theme_for(prs, template_entry)
    if template_entry is not None:
        _check_layout_provenance(prs, template_entry, findings)
        _check_density(prs, template_entry, findings)
        _check_slide_size(prs, template_entry, findings)
    _check_fonts(prs, theme, findings)
    _check_colors(prs, theme, findings)
    _check_footers(prs, findings)
    _check_master_elements_covered(prs, findings)
    _check_placeholder_bypass(prs, findings)
    _check_overflow(prs, findings)
    _check_offgrid(prs, findings)
    findings.sort(key=lambda f: (f["slide_index"] or 0, f["rule"]))
    return findings


# -- repair -----------------------------------------------------------------------


def plan_repairs(
    prs: PresentationType,
    findings: list[dict[str, Any]],
    template_entry: dict | None,
    strategy: str = "conservative",
) -> list[dict[str, Any]]:
    """Conservative strategy: only zero-risk fixes — re-link non-theme fonts
    to the theme, snap near-theme colors to their exact theme value."""
    if strategy != "conservative":
        raise ValueError(
            "Only the 'conservative' repair strategy is implemented; 'aggressive' "
            "(placeholder migration, reflow) arrives with apply_template."
        )
    theme = _theme_for(prs, template_entry)
    fixes = []
    for finding in findings:
        if not finding["auto_fixable"]:
            continue
        if finding["rule"] == "C02":
            fixes.append(
                {
                    "rule": "C02",
                    "slide_index": finding["slide_index"],
                    "shape_id": finding["shape_id"],
                    "action": "relink_font_to_theme",
                    "detail": finding["message"].split("'")[1],
                }
            )
        elif finding["rule"] == "C03":
            value = finding["message"].split(" ")[2]
            _, target_hex, _ = _nearest_theme_color(value, theme["color_scheme"])
            fixes.append(
                {
                    "rule": "C03",
                    "slide_index": finding["slide_index"],
                    "shape_id": finding["shape_id"],
                    "action": "snap_color_to_theme",
                    "from": value,
                    "to": target_hex,
                }
            )
    return fixes


def apply_repairs(prs: PresentationType, fixes: list[dict[str, Any]], template_entry: dict | None) -> int:
    from pptx.dml.color import RGBColor

    theme = _theme_for(prs, template_entry)
    fonts = theme["font_scheme"]
    applied = 0
    for fix in fixes:
        slide = prs.slides[fix["slide_index"] - 1]
        shape = next((s for s in slide.shapes if s.shape_id == fix["shape_id"]), None)
        if shape is None:
            continue
        if fix["action"] == "relink_font_to_theme":
            for run in _iter_runs(shape):
                if not _theme_font_ok(run.font.name, fonts):
                    run.font.name = None  # inherit the theme font again
            applied += 1
        elif fix["action"] == "snap_color_to_theme":
            for run in _iter_runs(shape):
                color = run.font.color
                try:
                    is_rgb = color.type == MSO_COLOR_TYPE.RGB
                except (AttributeError, ValueError):
                    continue
                if is_rgb and f"#{color.rgb}".upper() == fix["from"].upper():
                    color.rgb = RGBColor.from_string(fix["to"].lstrip("#"))
            applied += 1
    return applied

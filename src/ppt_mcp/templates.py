"""Template Intelligence Layer — registry and parser (DESIGN.md §5).

A registered template is parsed once into a registry entry (§4.4): masters,
layouts with placeholder schemas, theme color/font schemes, inferred intent
tags, and coarse capacity estimates. The source file is copied into the
store and treated as immutable from then on.
"""

from __future__ import annotations

import hashlib
import json
import shutil
import zipfile
from pathlib import Path
from typing import Any

from lxml import etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from .errors import PptMcpError
from .reader import emu_to_inches, load_presentation
from .store import Store
from .writer import placeholder_role

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"a": A_NS}

TEMPLATE_SUFFIXES = {".pptx", ".potx", ".thmx"}

# PresentationML main-part content types: template vs presentation.
CT_PRESENTATION = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
CT_TEMPLATE = "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"

# Layout-name keywords -> intent tags (multilingual; user-editable afterwards).
INTENT_KEYWORDS = {
    "title slide": "title_slide",
    "titelfolie": "title_slide",
    "section": "section",
    "divider": "section",
    "abschnitt": "section",
    "kapitel": "section",
    "agenda": "agenda",
    "inhalt": "agenda",
    "two content": "two_column",
    "zwei inhalte": "two_column",
    "two column": "two_column",
    "comparison": "comparison",
    "vergleich": "comparison",
    "picture": "image",
    "bild": "image",
    "image": "image",
    "photo": "image",
    "blank": "blank",
    "leer": "blank",
    "title only": "title_only",
    "nur titel": "title_only",
    "timeline": "timeline",
    "roadmap": "timeline",
    "zeitplan": "timeline",
    "table": "table",
    "tabelle": "table",
    "chart": "chart",
    "diagramm": "chart",
    "graph": "chart",
    "executive": "executive_summary",
    "summary": "executive_summary",
    "zusammenfassung": "executive_summary",
    "management summary": "executive_summary",
    "quote": "quote",
    "zitat": "quote",
    "team": "team",
    "org": "team",
    "appendix": "appendix",
    "anhang": "appendix",
    "matrix": "matrix",
    "swot": "matrix",
    "risk": "matrix",
    "risiko": "matrix",
    "dashboard": "dashboard",
    "kpi": "dashboard",
    "caption": "caption",
}

# Rough text-fit constants for capacity estimates (refined by rendering in M3).
_LINE_HEIGHT_IN = 0.42
_CHARS_PER_INCH = 11


def extract_theme_from_master(master) -> dict[str, Any]:
    """Color and font scheme straight from the theme part XML (python-pptx
    does not expose theme parsing — this is the lxml escape hatch)."""
    theme_part = master.part.part_related_by(RT.THEME)
    root = etree.fromstring(theme_part.blob)
    colors: dict[str, str] = {}
    scheme = root.find(f".//{{{A_NS}}}clrScheme")
    if scheme is not None:
        for child in scheme:
            name = etree.QName(child).localname
            srgb = child.find(f"{{{A_NS}}}srgbClr")
            sys = child.find(f"{{{A_NS}}}sysClr")
            if srgb is not None:
                colors[name] = f"#{srgb.get('val')}"
            elif sys is not None and sys.get("lastClr"):
                colors[name] = f"#{sys.get('lastClr')}"
    fonts: dict[str, str | None] = {}
    font_scheme = root.find(f".//{{{A_NS}}}fontScheme")
    if font_scheme is not None:
        for kind, tag in (("major", "majorFont"), ("minor", "minorFont")):
            latin = font_scheme.find(f"{{{A_NS}}}{tag}/{{{A_NS}}}latin")
            fonts[kind] = latin.get("typeface") if latin is not None else None
    return {"color_scheme": colors, "font_scheme": fonts}


def _layout_intent_tags(name: str, placeholders: list[dict[str, Any]]) -> list[str]:
    tags: set[str] = set()
    lowered = name.lower()
    for keyword, tag in INTENT_KEYWORDS.items():
        if keyword in lowered:
            tags.add(tag)
    roles = [ph["role"] for ph in placeholders]
    body_count = roles.count("body")
    if "picture" in roles:
        tags.add("image")
    if "chart" in roles:
        tags.add("chart")
    if "table" in roles:
        tags.add("table")
    if body_count >= 2:
        tags.add("two_column")
        tags.add("comparison")
    if "subtitle" in roles and body_count == 0:
        tags.add("title_slide")
    if roles == ["title"]:
        tags.add("title_only")
        tags.add("section")
    if not placeholders:
        tags.add("blank")
    return sorted(tags)


def _capacity(placeholders: list[dict[str, Any]]) -> dict[str, Any]:
    body = [ph for ph in placeholders if ph["role"] == "body"]
    if not body:
        return {"body_bullets": 0, "chars_per_line": 0, "estimated": True}
    heights = [ph["height_in"] for ph in body if ph.get("height_in")]
    widths = [ph["width_in"] for ph in body if ph.get("width_in")]
    return {
        "body_bullets": int(min(heights) / _LINE_HEIGHT_IN) if heights else 0,
        "chars_per_line": int(min(widths) * _CHARS_PER_INCH) if widths else 0,
        "estimated": True,
    }


def _placeholder_schema(layout) -> list[dict[str, Any]]:
    schema = []
    for shape in layout.placeholders:
        ph = shape.placeholder_format
        type_name = ph.type.name if ph.type is not None else "BODY"
        schema.append(
            {
                "idx": ph.idx,
                "type": type_name,
                "role": placeholder_role(shape),
                "name": shape.name,
                "left_in": emu_to_inches(shape.left),
                "top_in": emu_to_inches(shape.top),
                "width_in": emu_to_inches(shape.width),
                "height_in": emu_to_inches(shape.height),
            }
        )
    return schema


def parse_template(path: Path) -> dict[str, Any]:
    prs = load_presentation(path)
    masters = []
    layouts = []
    for m_i, master in enumerate(prs.slide_masters):
        master_id = f"master_{m_i + 1}"
        layout_ids = []
        for l_i, layout in enumerate(master.slide_layouts):
            layout_id = f"layout_{m_i + 1}_{l_i + 1}"
            schema = _placeholder_schema(layout)
            content_schema = [
                ph for ph in schema if ph["role"] not in ("footer", "date", "slide_number")
            ]
            layouts.append(
                {
                    "layout_id": layout_id,
                    "name": layout.name,
                    "master_id": master_id,
                    "placeholders": schema,
                    "intent_tags": _layout_intent_tags(layout.name, content_schema),
                    "capacity": _capacity(content_schema),
                }
            )
            layout_ids.append(layout_id)
        masters.append(
            {
                "master_id": master_id,
                "name": master.name or f"Master {m_i + 1}",
                "layout_ids": layout_ids,
            }
        )
    width, height = int(prs.slide_width), int(prs.slide_height)
    return {
        "slide_size": {
            "width_in": emu_to_inches(width),
            "height_in": emu_to_inches(height),
        },
        "masters": masters,
        "layouts": layouts,
        "theme": extract_theme_from_master(prs.slide_masters[0]),
        "example_slide_count": len(prs.slides),
    }


def materialize_as_pptx(source: Path, dest: Path) -> None:
    """Copy a template package as a presentation: byte copy, but with the main
    part's content type switched from template to presentation so PowerPoint
    treats the result as a normal deck (relevant for .potx sources)."""
    with zipfile.ZipFile(source) as zin, zipfile.ZipFile(
        dest, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(CT_TEMPLATE.encode(), CT_PRESENTATION.encode())
            zout.writestr(item, data)


class TemplateRegistry:
    def __init__(self, store: Store):
        self.store = store
        self.index_path = store.registry_dir / "templates.json"
        self._entries: dict[str, dict[str, Any]] = {}
        if self.index_path.is_file():
            self._entries = json.loads(self.index_path.read_text(encoding="utf-8"))

    def _save_index(self) -> None:
        self.index_path.write_text(
            json.dumps(self._entries, indent=2, ensure_ascii=False), encoding="utf-8"
        )

    def register(
        self,
        path: str | Path,
        name: str | None = None,
        version: str | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        source = Path(path).expanduser().resolve()
        if not source.is_file():
            raise PptMcpError(f"Template file not found: '{source}'.")
        if source.suffix.lower() not in TEMPLATE_SUFFIXES - {".thmx"}:
            raise PptMcpError(
                f"Unsupported template type '{source.suffix}'. Supported: .pptx, .potx "
                "(.thmx arrives with theme-only mode)."
            )
        digest = hashlib.sha256(source.read_bytes()).hexdigest()
        for existing in self._entries.values():
            if existing["sha256"] == digest:
                return {**existing, "already_registered": True}
        template_id = f"tpl_{digest[:8]}"
        template_dir = self.store.templates_dir / template_id
        template_dir.mkdir(parents=True, exist_ok=True)
        stored_source = template_dir / f"source{source.suffix.lower()}"
        shutil.copy2(source, stored_source)

        # python-pptx only opens presentation-typed packages, so .potx sources
        # are parsed (and later instantiated) through a content-type-patched copy
        materialized = template_dir / "materialized.pptx"
        materialize_as_pptx(stored_source, materialized)
        parsed = parse_template(materialized)
        (template_dir / "parsed.json").write_text(
            json.dumps(parsed, indent=2, ensure_ascii=False), encoding="utf-8"
        )
        entry = {
            "template_id": template_id,
            "name": name or source.stem,
            "version": version or "1.0",
            "source_path": str(stored_source),
            "materialized_path": str(materialized),
            "original_path": str(source),
            "sha256": digest,
            "builtin": False,
            "metadata": metadata or {},
            **parsed,
        }
        self._entries[template_id] = entry
        self._save_index()
        self.store.log_provenance(
            "ppt_register_template", template_id=template_id, source=source
        )
        return entry

    def get(self, template_id: str) -> dict[str, Any]:
        entry = self._entries.get(template_id)
        if entry is None:
            known = ", ".join(self._entries) or "none"
            raise PptMcpError(
                f"Unknown template_id '{template_id}'. Registered templates: {known}. "
                "Use ppt_register_template to add one."
            )
        return entry

    def list(self) -> list[dict[str, Any]]:
        return list(self._entries.values())

    def get_layout(self, template_id: str, layout_ref: str) -> dict[str, Any]:
        """Resolve by layout_id or layout name (case-insensitive)."""
        entry = self.get(template_id)
        for layout in entry["layouts"]:
            if layout["layout_id"] == layout_ref or layout["name"].lower() == layout_ref.lower():
                return layout
        names = ", ".join(f"{lo['layout_id']} ('{lo['name']}')" for lo in entry["layouts"])
        raise PptMcpError(
            f"Layout '{layout_ref}' not found in template {template_id}. "
            f"Available: {names}."
        )

    def update(self, template_id: str, patch: dict[str, Any]) -> dict[str, Any]:
        """Patch user-editable fields: name, version, metadata, and per-layout
        intent_tags ({'layout_intent_tags': {layout_id: [tags]}})."""
        entry = self.get(template_id)
        allowed = {"name", "version", "metadata"}
        unknown = set(patch) - allowed - {"layout_intent_tags"}
        if unknown:
            raise PptMcpError(
                f"Cannot patch {sorted(unknown)}. Editable: name, version, metadata, "
                "layout_intent_tags."
            )
        for key in allowed & set(patch):
            entry[key] = patch[key]
        for layout_id, tags in (patch.get("layout_intent_tags") or {}).items():
            layout = next(
                (lo for lo in entry["layouts"] if lo["layout_id"] == layout_id), None
            )
            if layout is None:
                raise PptMcpError(f"Unknown layout_id '{layout_id}' in {template_id}.")
            layout["intent_tags"] = sorted(set(tags))
        self._save_index()
        return entry

    def materialized_path(self, template_id: str) -> Path:
        """Presentation-typed copy of the template, safe for python-pptx."""
        return Path(self.get(template_id)["materialized_path"])

    def load_layout_object(self, template_id: str, layout_ref: str):
        """The live python-pptx layout object for mapping/recommendation."""
        layout_meta = self.get_layout(template_id, layout_ref)
        prs = load_presentation(self.materialized_path(template_id))
        m_index = int(layout_meta["master_id"].split("_")[1]) - 1
        l_index = int(layout_meta["layout_id"].split("_")[2]) - 1
        return prs.slide_masters[m_index].slide_layouts[l_index]

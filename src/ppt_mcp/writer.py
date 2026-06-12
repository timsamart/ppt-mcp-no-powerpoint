"""Write-side OOXML engine (DESIGN.md §10.3): placeholder-first authoring.

Every mutation has two phases so dry-run is free: `plan_*` computes a
structured change plan without touching the file; `apply_*` executes it.
Content goes into placeholders by semantic role — absolute positioning only
through explicitly gated freeform paths.
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide, SlideLayout
from pptx.util import Inches

from .errors import PptMcpError
from .models import (
    ChartSpec,
    ContentSpec,
    EditOp,
    ImageRef,
    Paragraph,
    Position,
    ShapeContent,
    TableSpec,
)
from .reader import ROLE_BY_PH_TYPE_NAME, get_slide_by_index

XL_TYPE_BY_NAME = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
}

# Placeholder roles that python-pptx does not clone onto new slides.
NON_CLONED_ROLES = {"footer", "date", "slide_number"}

_SHAPE_TAGS = tuple(
    qn(f"p:{tag}") for tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp")
)


# -- resolution -----------------------------------------------------------------


def resolve_layout(prs: PresentationType, layout_ref: str) -> SlideLayout:
    """Resolve a layout by name (case-insensitive)."""
    available = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            available.append(layout)
            if layout.name.lower() == layout_ref.lower():
                return layout
    names = ", ".join(f"'{lo.name}'" for lo in available)
    raise PptMcpError(
        f"Layout '{layout_ref}' not found. Available layouts: {names}. "
        "Use ppt_get_deck_overview to inspect masters and layouts."
    )


def placeholder_role(shape) -> str:
    ph_type = shape.placeholder_format.type
    type_name = ph_type.name if ph_type is not None else "BODY"
    return ROLE_BY_PH_TYPE_NAME.get(type_name, type_name.lower())


def resolve_shape(slide: Slide, shape_ref: str):
    """shape_ref forms (DESIGN.md §4.2): a shape_id ('7'), a placeholder idx
    ('idx:1'), or a semantic role ('title', 'body', 'picture', ...)."""
    ref = shape_ref.strip()
    if ref.isdigit():
        for shape in slide.shapes:
            if shape.shape_id == int(ref):
                return shape
        ids = ", ".join(str(s.shape_id) for s in slide.shapes)
        raise PptMcpError(
            f"No shape with shape_id {ref} on this slide. Shape ids present: {ids}. "
            "Use ppt_get_slide to inspect."
        )
    if ref.lower().startswith("idx:"):
        idx = int(ref[4:])
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == idx:
                return shape
        idxs = ", ".join(str(s.placeholder_format.idx) for s in slide.placeholders)
        raise PptMcpError(
            f"No placeholder with idx {idx} on this slide. Placeholder idxs: {idxs}."
        )
    matches = [s for s in slide.placeholders if placeholder_role(s) == ref.lower()]
    if len(matches) == 1:
        return matches[0]
    roles = ", ".join(sorted({placeholder_role(s) for s in slide.placeholders}))
    if not matches:
        raise PptMcpError(
            f"No placeholder with role '{ref}' on this slide. Roles present: "
            f"{roles or 'none'}. You can also address shapes by shape_id or 'idx:N'."
        )
    options = ", ".join(f"idx:{s.placeholder_format.idx}" for s in matches)
    raise PptMcpError(
        f"Role '{ref}' is ambiguous on this slide ({len(matches)} placeholders: "
        f"{options}). Address one of them by idx."
    )


# -- text frames ------------------------------------------------------------------


def set_text_frame(text_frame, paragraphs: list[Paragraph]) -> None:
    text_frame.clear()
    for i, para in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = para.text
        if para.level:
            p.level = para.level


def fill_table(graphic_frame, spec: TableSpec) -> None:
    table = graphic_frame.table
    data: list[list[str]] = ([spec.headers] if spec.headers else []) + spec.rows
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            table.cell(r, c).text = str(value)


def _table_dims(spec: TableSpec) -> tuple[int, int]:
    rows = len(spec.rows) + (1 if spec.headers else 0)
    cols = len(spec.headers) if spec.headers else len(spec.rows[0])
    return rows, cols


def chart_data(spec: ChartSpec) -> CategoryChartData:
    data = CategoryChartData()
    data.categories = spec.categories
    for series in spec.series:
        data.add_series(series.name, series.values)
    return data


# -- placeholder content insertion -------------------------------------------------


def insert_into_placeholder(shape, content: ShapeContent) -> dict[str, Any]:
    """Put one piece of content into one placeholder, dispatching on content
    kind. Returns a summary of what was done; raises with guidance when the
    placeholder cannot accept the content."""
    if content.text is not None:
        _require_text_frame(shape)
        set_text_frame(shape.text_frame, [Paragraph(text=content.text)])
        return {"set": "text"}
    if content.paragraphs is not None:
        _require_text_frame(shape)
        set_text_frame(shape.text_frame, content.paragraphs)
        return {"set": f"{len(content.paragraphs)} paragraph(s)"}
    if content.table is not None:
        if not hasattr(shape, "insert_table"):
            raise PptMcpError(
                f"Placeholder '{shape.name}' (role {placeholder_role(shape)}) cannot "
                "hold a table — only TABLE-type placeholders can. Use a layout with a "
                "table placeholder, or ppt_add_table with allow_freeform=true."
            )
        rows, cols = _table_dims(content.table)
        frame = shape.insert_table(rows=rows, cols=cols)
        fill_table(frame, content.table)
        return {"set": f"table {rows}x{cols}", "shape_id": frame.shape_id}
    if content.image is not None:
        if not hasattr(shape, "insert_picture"):
            raise PptMcpError(
                f"Placeholder '{shape.name}' (role {placeholder_role(shape)}) cannot "
                "hold a picture — only PICTURE-type placeholders can. Use "
                "ppt_add_image with allow_freeform=true for free placement."
            )
        picture = shape.insert_picture(_checked_image_path(content.image))
        if content.image.alt_text:
            picture._element._nvXxPr.cNvPr.set("descr", content.image.alt_text)
        return {"set": "picture", "shape_id": picture.shape_id}
    if content.chart is not None:
        if not hasattr(shape, "insert_chart"):
            raise PptMcpError(
                f"Placeholder '{shape.name}' (role {placeholder_role(shape)}) cannot "
                "hold a chart — only CHART-type placeholders can. Use ppt_add_chart "
                "with allow_freeform=true for free placement."
            )
        frame = shape.insert_chart(
            XL_TYPE_BY_NAME[content.chart.chart_type], chart_data(content.chart)
        )
        return {"set": f"{content.chart.chart_type} chart", "shape_id": frame.shape_id}
    raise PptMcpError("Empty content.")  # unreachable: ShapeContent validates


def _require_text_frame(shape) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise PptMcpError(
            f"Shape '{shape.name}' has no text frame; it cannot take text content."
        )


def _checked_image_path(image: ImageRef) -> str:
    path = Path(image.path).expanduser()
    if not path.is_file():
        raise PptMcpError(f"Image file not found: '{path}'.")
    return str(path)


# -- content mapping (DESIGN.md §5.4) -----------------------------------------------


def _layout_placeholder_inventory(layout: SlideLayout) -> list[dict[str, Any]]:
    # Capabilities come from the placeholder *type*: layout placeholders are
    # generic objects; the specialized insert_* classes only materialize on
    # slides, so hasattr checks would always be False here.
    inventory = []
    for shape in layout.placeholders:
        role = placeholder_role(shape)
        if role in NON_CLONED_ROLES:
            continue
        ph_type = shape.placeholder_format.type
        type_name = ph_type.name if ph_type is not None else "BODY"
        inventory.append(
            {
                "idx": shape.placeholder_format.idx,
                "role": role,
                "can_table": type_name == "TABLE",
                "can_picture": type_name in ("PICTURE", "CLIP_ART", "BITMAP"),
                "can_chart": type_name == "CHART",
            }
        )
    return inventory


def plan_content_mapping(layout: SlideLayout, spec: ContentSpec) -> dict[str, Any]:
    """Pure mapping: spec sections -> layout placeholders. Sections with no
    suitable placeholder land in unplaced_content with advice."""
    available = _layout_placeholder_inventory(layout)
    used: set[int] = set()
    placements: list[dict[str, Any]] = []
    unplaced: list[str] = []

    def take(predicate, section: str, summary: str) -> None:
        for ph in available:
            if ph["idx"] not in used and predicate(ph):
                used.add(ph["idx"])
                placements.append(
                    {
                        "section": section,
                        "placeholder_idx": ph["idx"],
                        "role": ph["role"],
                        "summary": summary,
                    }
                )
                return
        unplaced.append(section)

    if spec.title is not None:
        take(lambda ph: ph["role"] == "title", "title", spec.title)
    if spec.subtitle is not None:
        take(lambda ph: ph["role"] == "subtitle", "subtitle", spec.subtitle)
    if spec.body is not None:
        take(
            lambda ph: ph["role"] == "body",
            "body",
            f"{len(spec.body)} paragraph(s)",
        )
    if spec.table is not None:
        rows, cols = _table_dims(spec.table)
        take(lambda ph: ph["can_table"], "table", f"{rows}x{cols} table")
    if spec.image is not None:
        take(lambda ph: ph["can_picture"], "image", spec.image.path)
    if spec.chart is not None:
        take(
            lambda ph: ph["can_chart"],
            "chart",
            f"{spec.chart.chart_type} chart, {len(spec.chart.series)} series",
        )

    unplaced_messages = [
        f"{section}: no free placeholder in layout '{layout.name}' can hold it. "
        "Pick a different layout (ppt_get_deck_overview lists them) or place it "
        "explicitly with the dedicated add tool."
        for section in unplaced
    ]
    return {
        "layout": layout.name,
        "placements": placements,
        "unplaced_content": unplaced_messages,
        "notes": spec.notes is not None,
    }


def _spec_content_for_section(spec: ContentSpec, section: str) -> ShapeContent:
    if section == "title":
        return ShapeContent(text=spec.title)
    if section == "subtitle":
        return ShapeContent(text=spec.subtitle)
    if section == "body":
        return ShapeContent(paragraphs=spec.body)
    if section == "table":
        return ShapeContent(table=spec.table)
    if section == "image":
        return ShapeContent(image=spec.image)
    if section == "chart":
        return ShapeContent(chart=spec.chart)
    raise PptMcpError(f"Unknown content section '{section}'.")


def apply_add_slide(
    prs: PresentationType,
    layout: SlideLayout,
    spec: ContentSpec,
    plan: dict[str, Any],
    position: int | None = None,
) -> int:
    """Create the slide and execute a plan from plan_content_mapping.
    Returns the new slide's 1-based index."""
    slide = prs.slides.add_slide(layout)
    for placement in plan["placements"]:
        target = next(
            ph
            for ph in slide.placeholders
            if ph.placeholder_format.idx == placement["placeholder_idx"]
        )
        insert_into_placeholder(target, _spec_content_for_section(spec, placement["section"]))
    if spec.notes is not None:
        slide.notes_slide.notes_text_frame.text = spec.notes
    new_index = len(prs.slides)
    if position is not None and position != new_index:
        move_slide(prs, new_index, position)
        new_index = position
    return new_index


# -- text edits -----------------------------------------------------------------


def apply_edit_ops(shape, ops: list[EditOp]) -> dict[str, Any]:
    _require_text_frame(shape)
    text_frame = shape.text_frame
    replacements = 0
    formatting_warnings = 0
    for op in ops:
        if op.op == "replace_text":
            for paragraph in text_frame.paragraphs:
                haystack = paragraph.text if op.match_case else paragraph.text.lower()
                needle = op.find if op.match_case else op.find.lower()
                if needle not in haystack:
                    continue
                hit_in_single_run = False
                for run in paragraph.runs:
                    run_haystack = run.text if op.match_case else run.text.lower()
                    if needle in run_haystack:
                        run.text = _ci_replace(run.text, op.find, op.replace, op.match_case)
                        hit_in_single_run = True
                        replacements += 1
                if not hit_in_single_run:
                    # match spans runs: rewrite the paragraph, losing run-level
                    # formatting — counted so the caller can see it happened
                    paragraph.text = _ci_replace(
                        paragraph.text, op.find, op.replace, op.match_case
                    )
                    replacements += 1
                    formatting_warnings += 1
        elif op.op == "set_paragraphs":
            set_text_frame(text_frame, op.paragraphs)
        elif op.op == "append_paragraph":
            p = text_frame.add_paragraph()
            p.text = op.text
            if op.level:
                p.level = op.level
    result: dict[str, Any] = {"replacements": replacements}
    if formatting_warnings:
        result["warnings"] = [
            f"{formatting_warnings} replacement(s) spanned formatting runs; "
            "run-level formatting was flattened for those paragraphs."
        ]
    return result


def _ci_replace(text: str, find: str, replace: str, match_case: bool) -> str:
    if match_case:
        return text.replace(find, replace)
    out, low, needle = [], text.lower(), find.lower()
    i = 0
    while True:
        j = low.find(needle, i)
        if j < 0:
            out.append(text[i:])
            return "".join(out)
        out.append(text[i:j])
        out.append(replace)
        i = j + len(find)


# -- freeform additions (gated; DESIGN.md §10.3) -------------------------------------


def add_freeform_table(slide: Slide, spec: TableSpec, pos: Position):
    rows, cols = _table_dims(spec)
    frame = slide.shapes.add_table(
        rows, cols, Inches(pos.left), Inches(pos.top), Inches(pos.width), Inches(pos.height)
    )
    fill_table(frame, spec)
    return frame


def add_freeform_image(slide: Slide, image: ImageRef, pos: Position):
    picture = slide.shapes.add_picture(
        _checked_image_path(image),
        Inches(pos.left),
        Inches(pos.top),
        width=Inches(pos.width),
        height=Inches(pos.height),
    )
    if image.alt_text:
        picture._element._nvXxPr.cNvPr.set("descr", image.alt_text)
    return picture


def add_freeform_chart(slide: Slide, spec: ChartSpec, pos: Position):
    return slide.shapes.add_chart(
        XL_TYPE_BY_NAME[spec.chart_type],
        Inches(pos.left),
        Inches(pos.top),
        Inches(pos.width),
        Inches(pos.height),
        chart_data(spec),
    )


FREEFORM_WARNING = (
    "Freeform placement bypasses the layout's placeholder grammar (compliance "
    "rule C06). Prefer a layout with a suitable placeholder."
)


# -- slide operations -----------------------------------------------------------


def delete_slide(prs: PresentationType, slide_index: int) -> None:
    get_slide_by_index(prs, slide_index)  # bounds check
    sld_id_lst = prs.slides._sldIdLst
    element = list(sld_id_lst)[slide_index - 1]
    prs.part.drop_rel(element.get(qn("r:id")))
    sld_id_lst.remove(element)


def move_slide(prs: PresentationType, slide_index: int, to_position: int) -> None:
    count = len(prs.slides)
    get_slide_by_index(prs, slide_index)
    if not 1 <= to_position <= count:
        raise PptMcpError(f"to_position {to_position} out of range 1..{count}.")
    sld_id_lst = prs.slides._sldIdLst
    element = list(sld_id_lst)[slide_index - 1]
    sld_id_lst.remove(element)
    sld_id_lst.insert(to_position - 1, element)


def clone_shape_elements(
    source: Slide, dest: Slide, elements, reassign_ids: bool = False
) -> list[int]:
    """Deep-copy selected shape elements from one slide to another, re-creating
    the underlying part relationships (images, charts, hyperlinks) and
    remapping r:embed/r:link/r:id inside the copied XML. With reassign_ids,
    cloned shapes get fresh ids unique on the destination slide. Returns the
    cloned shapes' ids."""
    rid_map: dict[str, str] = {}
    for rid, rel in source.part.rels.items():
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        if rel.is_external:
            rid_map[rid] = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            rid_map[rid] = dest.part.relate_to(rel.target_part, rel.reltype)

    rel_attrs = (qn("r:embed"), qn("r:link"), qn("r:id"))
    next_id = max((s.shape_id for s in dest.shapes), default=1) + 1
    new_ids: list[int] = []
    for child in elements:
        clone = deepcopy(child)
        for element in clone.iter():
            for attr in rel_attrs:
                value = element.get(attr)
                if value and value in rid_map:
                    element.set(attr, rid_map[value])
        cnvpr = clone.find(f".//{qn('p:cNvPr')}")
        if reassign_ids and cnvpr is not None:
            cnvpr.set("id", str(next_id))
            next_id += 1
        dest.shapes._spTree.append(clone)
        if cnvpr is not None:
            new_ids.append(int(cnvpr.get("id")))
    return new_ids


def duplicate_slide(prs: PresentationType, slide_index: int) -> int:
    """Deep-copy a slide's shape tree and relationships onto a fresh slide.
    The copy is placed directly after the source. Returns its 1-based index."""
    source = get_slide_by_index(prs, slide_index)
    dest = prs.slides.add_slide(source.slide_layout)

    for element in list(dest.shapes._spTree):
        if element.tag in _SHAPE_TAGS:
            dest.shapes._spTree.remove(element)

    clone_shape_elements(
        source, dest, [c for c in source.shapes._spTree if c.tag in _SHAPE_TAGS]
    )

    if source.has_notes_slide:
        dest.notes_slide.notes_text_frame.text = (
            source.notes_slide.notes_text_frame.text
        )

    new_index = slide_index + 1
    move_slide(prs, len(prs.slides), new_index)
    return new_index

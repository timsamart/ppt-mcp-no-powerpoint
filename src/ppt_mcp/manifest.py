"""Image-placeholder manifest embedded in the deck (DESIGN.md §4.6).

Long generation prompts must not pollute visible slides, but they must
travel *with* the file. They live in a custom XML part
(`/customXml/pptmcp-manifest.xml`, namespace urn:ppt-mcp:manifest:v1) that
python-pptx round-trips untouched and PowerPoint preserves as an unknown
customXml part. Slides carry only a short label / notes pointer.
"""

from __future__ import annotations

import json
from typing import Any

from lxml import etree
from pptx.presentation import Presentation as PresentationType

MANIFEST_PARTNAME = "/customXml/pptmcp-manifest.xml"
MANIFEST_NS = "urn:ppt-mcp:manifest:v1"
RT_CUSTOM_XML = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/customXml"
)

_EMPTY: dict[str, Any] = {"version": 1, "image_placeholders": []}


def _find_part(prs: PresentationType):
    for rel in prs.part.rels.values():
        if (
            rel.reltype == RT_CUSTOM_XML
            and not rel.is_external
            and str(rel.target_part.partname) == MANIFEST_PARTNAME
        ):
            return rel.target_part
    return None


def load(prs: PresentationType) -> dict[str, Any]:
    part = _find_part(prs)
    if part is None:
        return json.loads(json.dumps(_EMPTY))  # fresh copy
    root = etree.fromstring(part.blob)
    payload = root.findtext(f"{{{MANIFEST_NS}}}json")
    return json.loads(payload) if payload else json.loads(json.dumps(_EMPTY))


def save(prs: PresentationType, data: dict[str, Any]) -> None:
    root = etree.Element(f"{{{MANIFEST_NS}}}manifest", nsmap={"m": MANIFEST_NS})
    holder = etree.SubElement(root, f"{{{MANIFEST_NS}}}json")
    holder.text = etree.CDATA(json.dumps(data, ensure_ascii=False))
    blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    part = _find_part(prs)
    if part is not None:
        part._blob = blob
        return
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    part = Part(
        partname=PackURI(MANIFEST_PARTNAME),
        content_type="application/xml",
        package=prs.part.package,
        blob=blob,
    )
    prs.part.relate_to(part, RT_CUSTOM_XML)


def find_record(
    data: dict[str, Any], slide_index: int, shape_id: int
) -> dict[str, Any] | None:
    for record in data["image_placeholders"]:
        if record["slide_index"] == slide_index and record["shape_id"] == shape_id:
            return record
    return None


def relocate_records(data: dict[str, Any], prs: PresentationType) -> None:
    """Slides move; shape_ids don't. Re-anchor records whose slide_index is
    stale by scanning for their shape_id."""
    for record in data["image_placeholders"]:
        index = record["slide_index"]
        if 1 <= index <= len(prs.slides) and any(
            s.shape_id == record["shape_id"] for s in prs.slides[index - 1].shapes
        ):
            continue
        for i, slide in enumerate(prs.slides):
            if any(s.shape_id == record["shape_id"] for s in slide.shapes):
                record["slide_index"] = i + 1
                break

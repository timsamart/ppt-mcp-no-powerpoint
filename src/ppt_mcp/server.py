"""ppt_mcp — FastMCP stdio server (DESIGN.md §10, M0 surface).

Deck lifecycle + reading/inspection tools. All logging goes to stderr
(stdio transport owns stdout). The server makes no network calls, ever.
"""

from __future__ import annotations

import logging
import sys
from pathlib import Path
from typing import Annotated, Any, Literal

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pydantic import Field

from . import compliance
from . import format as fmt
from . import manifest as deck_manifest
from . import reader, recommend, retarget, writer
from .errors import PptMcpError
from .render import RenderService, diff_images
from .styles import StyleProfileRegistry
from .templates import TemplateRegistry, extract_theme_from_master
from .models import (
    ChartSpec,
    ContentSpec,
    EditOp,
    ImagePlaceholderSpec,
    ImageRef,
    Position,
    ShapeContent,
    TableSpec,
)
from .sessions import SessionManager
from .store import Store

logging.basicConfig(stream=sys.stderr, level=logging.INFO)
log = logging.getLogger("ppt_mcp")

mcp = FastMCP("ppt_mcp")

store = Store()
sessions = SessionManager(store)
registry = TemplateRegistry(store)
renderer = RenderService(store)
styles = StyleProfileRegistry(store)

READ_ONLY = ToolAnnotations(readOnlyHint=True, openWorldHint=False)
MUTATING = ToolAnnotations(readOnlyHint=False, destructiveHint=False, openWorldHint=False)
DESTRUCTIVE = ToolAnnotations(readOnlyHint=False, destructiveHint=True, openWorldHint=False)

ResponseFormat = Literal["markdown", "json"]


def _prs(deck_id: str):
    return reader.load_presentation(sessions.get(deck_id).working_path)


def _commit(deck_id: str, prs, tool: str, **provenance: object) -> None:
    """Snapshot-then-write for every mutation: snapshot the pre-state, save
    the changed presentation to the working copy, log provenance."""
    sessions.snapshot(deck_id)
    prs.save(str(sessions.get(deck_id).working_path))
    store.log_provenance(tool, deck_id=deck_id, **provenance)


# -- deck lifecycle (§10.1) ----------------------------------------------------


@mcp.tool(annotations=MUTATING)
def ppt_open_deck(
    path: Annotated[str, Field(description="Absolute path to a .pptx or .potx file")],
) -> dict[str, Any]:
    """Open an existing PowerPoint deck for reading/editing. The file is copied
    into an isolated session working copy; the source is not touched until
    ppt_save_deck. Returns the deck_id handle plus a deck overview."""
    session = sessions.open_deck(path)
    overview = reader.deck_overview(reader.load_presentation(session.working_path))
    return {
        "deck_id": session.deck_id,
        "source_path": str(session.source_path),
        "overview": overview,
    }


@mcp.tool(annotations=MUTATING)
def ppt_create_deck(
    template_id: Annotated[
        str | None,
        Field(description="Registered template to build on (ppt_list_templates). "
                          "Inherits its masters, layouts, and theme."),
    ] = None,
    include_example_slides: Annotated[
        bool, Field(description="Keep the template's example slides instead of starting empty")
    ] = False,
) -> dict[str, Any]:
    """Create a new deck — template-first: pass a registered template_id to
    inherit the corporate design system. Without one you get a bare default
    deck. Save later with ppt_save_deck(path=...)."""
    if template_id is None:
        session = sessions.create_deck()
        return {"deck_id": session.deck_id, "source_path": None, "template_id": None}
    entry = registry.get(template_id)
    session = sessions.create_deck(template_source=registry.materialized_path(template_id))
    prs = reader.load_presentation(session.working_path)
    removed = 0
    if not include_example_slides:
        for index in range(len(prs.slides), 0, -1):
            writer.delete_slide(prs, index)
            removed += 1
        prs.save(str(session.working_path))
    return {
        "deck_id": session.deck_id,
        "template_id": template_id,
        "template_name": entry["name"],
        "example_slides_removed": removed,
        "layouts": [lo["name"] for lo in entry["layouts"]],
    }


@mcp.tool(annotations=MUTATING)
def ppt_save_deck(
    deck_id: str,
    path: Annotated[
        str | None,
        Field(description="Target path. Defaults to the path the deck was opened from."),
    ] = None,
) -> dict[str, Any]:
    """Write the deck's working copy to disk. Without a path, saves back to the
    source file the deck was opened from."""
    target = sessions.save_deck(deck_id, path)
    return {"deck_id": deck_id, "saved_to": str(target)}


@mcp.tool(annotations=DESTRUCTIVE)
def ppt_close_deck(deck_id: str) -> dict[str, Any]:
    """Close a deck session and delete its working copy and snapshots.
    Unsaved changes are lost — call ppt_save_deck first if needed."""
    sessions.close_deck(deck_id)
    return {"deck_id": deck_id, "closed": True}


@mcp.tool(annotations=READ_ONLY)
def ppt_list_decks() -> dict[str, Any]:
    """List open deck sessions with their ids and source paths."""
    return {
        "decks": [
            {
                "deck_id": s.deck_id,
                "source_path": str(s.source_path) if s.source_path else None,
                "opened_at": s.opened_at,
                "undo_steps_available": s.snapshot_count,
            }
            for s in sessions.list_sessions()
        ]
    }


@mcp.tool(annotations=MUTATING)
def ppt_undo(
    deck_id: str,
    steps: Annotated[int, Field(ge=1, description="Number of mutations to roll back")] = 1,
) -> dict[str, Any]:
    """Undo the last N mutations on a deck by restoring a pre-mutation snapshot."""
    undone = sessions.undo(deck_id, steps)
    remaining = sessions.get(deck_id).snapshot_count
    if undone == 0:
        return {"deck_id": deck_id, "undone_steps": 0, "note": "Nothing to undo."}
    return {"deck_id": deck_id, "undone_steps": undone, "undo_steps_remaining": remaining}


# -- reading & inspection (§10.2) -----------------------------------------------


@mcp.tool(annotations=READ_ONLY)
def ppt_get_deck_overview(
    deck_id: str,
    response_format: ResponseFormat = "markdown",
) -> str | dict[str, Any]:
    """Deck summary: slide size/aspect ratio, masters with their layouts, and a
    one-line-per-slide listing (index, layout, title)."""
    data = reader.deck_overview(_prs(deck_id))
    if response_format == "json":
        return {"deck_id": deck_id, **data}
    return fmt.overview_markdown(deck_id, data)


@mcp.tool(annotations=READ_ONLY)
def ppt_get_slide(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1, description="1-based, as in the PowerPoint UI")],
    response_format: ResponseFormat = "markdown",
) -> str | dict[str, Any]:
    """Full inventory of one slide: every shape with its shape_id, placeholder
    role/idx, geometry (inches + EMU, layout-inherited where applicable), text
    with indent levels, tables/charts, and speaker notes."""
    data = reader.slide_detail(_prs(deck_id), slide_index)
    if response_format == "json":
        return {"deck_id": deck_id, **data}
    return fmt.slide_markdown(data)


@mcp.tool(annotations=READ_ONLY)
def ppt_search_deck(
    deck_id: str,
    query: Annotated[str, Field(min_length=1, description="Case-insensitive substring")],
    limit: Annotated[int, Field(ge=1, le=200)] = 25,
    offset: Annotated[int, Field(ge=0)] = 0,
    response_format: ResponseFormat = "markdown",
) -> str | dict[str, Any]:
    """Search all slide text and speaker notes. Returns hits with slide index
    and shape_id so results can be addressed directly by follow-up calls."""
    data = reader.search_deck(_prs(deck_id), query, limit=limit, offset=offset)
    if response_format == "json":
        return {"deck_id": deck_id, **data}
    return fmt.search_markdown(data)


# -- template registry & intelligence (§10.4) -------------------------------------


@mcp.tool(annotations=MUTATING)
def ppt_register_template(
    path: Annotated[str, Field(description="Absolute path to a .potx or template .pptx")],
    name: str | None = None,
    version: str | None = None,
    metadata: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Register a corporate template: the file is copied into the local store
    (immutable from then on) and parsed into a design-system profile — masters,
    layouts with placeholder schemas, theme colors/fonts, inferred intent tags."""
    entry = registry.register(path, name=name, version=version, metadata=metadata)
    return {
        "already_registered": entry.get("already_registered", False),
        "template_id": entry["template_id"],
        "name": entry["name"],
        "version": entry["version"],
        "masters": len(entry["masters"]),
        "layouts": len(entry["layouts"]),
        "example_slide_count": entry["example_slide_count"],
        "theme_accents": {
            k: v for k, v in entry["theme"]["color_scheme"].items() if k.startswith("accent")
        },
    }


@mcp.tool(annotations=READ_ONLY)
def ppt_list_templates() -> dict[str, Any]:
    """List registered templates."""
    return {
        "templates": [
            {
                "template_id": e["template_id"],
                "name": e["name"],
                "version": e["version"],
                "layouts": len(e["layouts"]),
                "builtin": e.get("builtin", False),
            }
            for e in registry.list()
        ]
    }


@mcp.tool(annotations=READ_ONLY)
def ppt_inspect_template(
    template_id: str,
    response_format: ResponseFormat = "markdown",
) -> str | dict[str, Any]:
    """Template design system: masters, layouts with placeholder roles and
    intent tags, theme color/font schemes."""
    entry = registry.get(template_id)
    if response_format == "json":
        return entry
    return fmt.template_markdown(entry)


@mcp.tool(annotations=READ_ONLY)
def ppt_inspect_layout(
    template_id: str,
    layout: Annotated[str, Field(description="layout_id or layout name")],
    response_format: ResponseFormat = "markdown",
) -> str | dict[str, Any]:
    """Full placeholder schema of one layout: idx, role, type, geometry, and
    estimated text capacity."""
    layout_meta = registry.get_layout(template_id, layout)
    if response_format == "json":
        return {"template_id": template_id, **layout_meta}
    return fmt.layout_markdown(template_id, layout_meta)


@mcp.tool(annotations=MUTATING)
def ppt_update_template(
    template_id: str,
    patch: Annotated[
        dict[str, Any],
        Field(description="Editable: name, version, metadata, layout_intent_tags "
                          "({layout_id: [tags]})"),
    ],
) -> dict[str, Any]:
    """Curate a registered template: rename, set version/metadata, or correct
    a layout's inferred intent tags."""
    entry = registry.update(template_id, patch)
    return {"template_id": template_id, "name": entry["name"], "updated": sorted(patch)}


@mcp.tool(annotations=READ_ONLY)
def ppt_extract_theme(
    template_id: str | None = None,
    deck_id: str | None = None,
) -> dict[str, Any]:
    """Theme color scheme and font scheme — from a registered template or from
    an open deck. Pass exactly one of template_id / deck_id."""
    if (template_id is None) == (deck_id is None):
        raise PptMcpError("Pass exactly one of template_id or deck_id.")
    if template_id is not None:
        return {"template_id": template_id, **registry.get(template_id)["theme"]}
    prs = _prs(deck_id)
    return {"deck_id": deck_id, **extract_theme_from_master(prs.slide_masters[0])}


@mcp.tool(annotations=READ_ONLY)
def ppt_recommend_layout(
    template_id: str,
    slide_intent: Annotated[str, Field(description="What the slide is for, e.g. 'risk overview'")],
    content: ContentSpec,
    top_n: Annotated[int, Field(ge=1, le=10)] = 3,
) -> dict[str, Any]:
    """Rank the template's layouts for a slide intent + content shape.
    Deterministic scoring with reasons — treat it as advice, not verdict."""
    entry = registry.get(template_id)
    recommendations = recommend.recommend_layouts(
        entry["layouts"], slide_intent, content, top_n=top_n
    )
    return {"template_id": template_id, "recommendations": recommendations}


@mcp.tool(annotations=READ_ONLY)
def ppt_map_content_to_placeholders(
    template_id: str,
    layout: Annotated[str, Field(description="layout_id or layout name")],
    content: ContentSpec,
) -> dict[str, Any]:
    """Preview how semantic content would map onto a layout's placeholders —
    the same mapping ppt_add_slide will apply, without touching any deck."""
    layout_obj = registry.load_layout_object(template_id, layout)
    return writer.plan_content_mapping(layout_obj, content)


# -- authoring (§10.3) -----------------------------------------------------------


@mcp.tool(annotations=MUTATING)
def ppt_add_slide(
    deck_id: str,
    layout: Annotated[str, Field(description="Layout name; ppt_get_deck_overview lists them")],
    content: ContentSpec,
    position: Annotated[
        int | None, Field(ge=1, description="1-based target position; default: append")
    ] = None,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Create a slide from a layout and fill its placeholders from semantic
    content (title/subtitle/body/table/image/chart/notes). Content that fits no
    placeholder is reported in unplaced_content — never silently dropped. Use
    dry_run=true to preview the mapping plan."""
    prs = _prs(deck_id)
    layout_obj = writer.resolve_layout(prs, layout)
    plan = writer.plan_content_mapping(layout_obj, content)
    if position is not None and position > len(prs.slides) + 1:
        raise PptMcpError(
            f"position {position} out of range; deck has {len(prs.slides)} slide(s)."
        )
    if dry_run:
        return {"applied": False, "plan": plan}
    new_index = writer.apply_add_slide(prs, layout_obj, content, plan, position)
    _commit(deck_id, prs, "ppt_add_slide", layout=layout_obj.name, slide_index=new_index)
    return {"applied": True, "slide_index": new_index, "plan": plan}


@mcp.tool(annotations=MUTATING)
def ppt_set_placeholder_content(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    shape_ref: Annotated[
        str, Field(description="Role ('title', 'body', ...), 'idx:N', or a shape_id")
    ],
    content: ShapeContent,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Replace the content of one placeholder, addressed by semantic role,
    placeholder idx, or shape_id. Exactly one content kind must be given."""
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)
    shape = writer.resolve_shape(slide, shape_ref)
    plan = {
        "slide_index": slide_index,
        "shape_id": shape.shape_id,
        "shape_name": shape.name,
        "content_kind": next(k for k, v in content.__dict__.items() if v is not None),
    }
    if dry_run:
        return {"applied": False, "plan": plan}
    result = writer.insert_into_placeholder(shape, content)
    _commit(deck_id, prs, "ppt_set_placeholder_content", slide_index=slide_index)
    return {"applied": True, "plan": plan, **result}


@mcp.tool(annotations=MUTATING)
def ppt_edit_text(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    shape_ref: str,
    ops: Annotated[
        list[EditOp],
        Field(description="Sequence of replace_text / set_paragraphs / append_paragraph ops"),
    ],
    dry_run: bool = False,
) -> dict[str, Any]:
    """Edit text of one shape with targeted operations. replace_text keeps
    run-level formatting when the match lies within a single run."""
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)
    shape = writer.resolve_shape(slide, shape_ref)
    plan = {
        "slide_index": slide_index,
        "shape_id": shape.shape_id,
        "ops": [op.op for op in ops],
    }
    if dry_run:
        return {"applied": False, "plan": plan}
    result = writer.apply_edit_ops(shape, ops)
    _commit(deck_id, prs, "ppt_edit_text", slide_index=slide_index)
    return {"applied": True, "plan": plan, **result}


@mcp.tool(annotations=MUTATING)
def ppt_set_notes(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    notes: str,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Set the speaker notes of a slide (replaces existing notes)."""
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)
    if dry_run:
        return {"applied": False, "plan": {"slide_index": slide_index, "set": "notes"}}
    slide.notes_slide.notes_text_frame.text = notes
    _commit(deck_id, prs, "ppt_set_notes", slide_index=slide_index)
    return {"applied": True, "slide_index": slide_index}


def _placed_add(
    deck_id: str,
    slide_index: int,
    content: ShapeContent,
    shape_ref: str | None,
    position: Position | None,
    allow_freeform: bool,
    dry_run: bool,
    freeform_apply,
    tool: str,
) -> dict[str, Any]:
    """Shared placeholder-first/freeform-gated flow for table/image/chart."""
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)
    kind = next(k for k, v in content.__dict__.items() if v is not None)
    if shape_ref is not None:
        shape = writer.resolve_shape(slide, shape_ref)
        plan = {"slide_index": slide_index, "target": f"placeholder shape_id {shape.shape_id}", "kind": kind}
        if dry_run:
            return {"applied": False, "plan": plan}
        result = writer.insert_into_placeholder(shape, content)
        _commit(deck_id, prs, tool, slide_index=slide_index, mode="placeholder")
        return {"applied": True, "plan": plan, **result}
    if not allow_freeform or position is None:
        raise PptMcpError(
            f"No shape_ref given. Either target a suitable placeholder (see "
            f"ppt_get_slide), or pass allow_freeform=true together with an explicit "
            f"position to place the {kind} freely. {writer.FREEFORM_WARNING}"
        )
    plan = {"slide_index": slide_index, "target": "freeform", "kind": kind,
            "warnings": [writer.FREEFORM_WARNING]}
    if dry_run:
        return {"applied": False, "plan": plan}
    new_shape = freeform_apply(slide, position)
    _commit(deck_id, prs, tool, slide_index=slide_index, mode="freeform")
    return {"applied": True, "plan": plan, "shape_id": new_shape.shape_id,
            "warnings": [writer.FREEFORM_WARNING]}


@mcp.tool(annotations=MUTATING)
def ppt_add_table(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    table: TableSpec,
    shape_ref: Annotated[str | None, Field(description="Target placeholder (preferred)")] = None,
    position: Position | None = None,
    allow_freeform: bool = False,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Add a table — into a table placeholder when shape_ref is given
    (preferred), or freeform only with allow_freeform=true and a position."""
    return _placed_add(
        deck_id, slide_index, ShapeContent(table=table), shape_ref, position,
        allow_freeform, dry_run,
        lambda slide, pos: writer.add_freeform_table(slide, table, pos),
        "ppt_add_table",
    )


@mcp.tool(annotations=MUTATING)
def ppt_add_image(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    image: ImageRef,
    shape_ref: Annotated[str | None, Field(description="Target picture placeholder (preferred)")] = None,
    position: Position | None = None,
    allow_freeform: bool = False,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Add a local image — into a picture placeholder when shape_ref is given
    (preferred), or freeform only with allow_freeform=true and a position."""
    return _placed_add(
        deck_id, slide_index, ShapeContent(image=image), shape_ref, position,
        allow_freeform, dry_run,
        lambda slide, pos: writer.add_freeform_image(slide, image, pos),
        "ppt_add_image",
    )


@mcp.tool(annotations=MUTATING)
def ppt_add_chart(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    chart: ChartSpec,
    shape_ref: Annotated[str | None, Field(description="Target chart placeholder (preferred)")] = None,
    position: Position | None = None,
    allow_freeform: bool = False,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Add a native chart (column/bar/line/pie/doughnut/area) — into a chart
    placeholder when shape_ref is given, or freeform with allow_freeform=true."""
    return _placed_add(
        deck_id, slide_index, ShapeContent(chart=chart), shape_ref, position,
        allow_freeform, dry_run,
        lambda slide, pos: writer.add_freeform_chart(slide, chart, pos),
        "ppt_add_chart",
    )


@mcp.tool(annotations=DESTRUCTIVE)
def ppt_delete_slide(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    dry_run: bool = False,
) -> dict[str, Any]:
    """Delete a slide. Recoverable via ppt_undo while the session is open."""
    prs = _prs(deck_id)
    title = reader.slide_title(reader.get_slide_by_index(prs, slide_index))
    plan = {"delete": {"slide_index": slide_index, "title": title}}
    if dry_run:
        return {"applied": False, "plan": plan}
    writer.delete_slide(prs, slide_index)
    _commit(deck_id, prs, "ppt_delete_slide", slide_index=slide_index)
    return {"applied": True, "plan": plan, "slide_count": len(prs.slides)}


@mcp.tool(annotations=MUTATING)
def ppt_move_slide(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    to_position: Annotated[int, Field(ge=1)],
    dry_run: bool = False,
) -> dict[str, Any]:
    """Move a slide to a new 1-based position."""
    prs = _prs(deck_id)
    plan = {"move": {"from": slide_index, "to": to_position}}
    if dry_run:
        return {"applied": False, "plan": plan}
    writer.move_slide(prs, slide_index, to_position)
    _commit(deck_id, prs, "ppt_move_slide", slide_index=slide_index, to=to_position)
    return {"applied": True, "plan": plan}


@mcp.tool(annotations=MUTATING)
def ppt_duplicate_slide(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    dry_run: bool = False,
) -> dict[str, Any]:
    """Duplicate a slide (shapes, images, notes); the copy lands right after
    the source."""
    prs = _prs(deck_id)
    title = reader.slide_title(reader.get_slide_by_index(prs, slide_index))
    plan = {"duplicate": {"slide_index": slide_index, "title": title}}
    if dry_run:
        return {"applied": False, "plan": plan}
    new_index = writer.duplicate_slide(prs, slide_index)
    _commit(deck_id, prs, "ppt_duplicate_slide", slide_index=slide_index)
    return {"applied": True, "plan": plan, "new_slide_index": new_index}


# -- brand style profiles & image placeholders (§10.6) ------------------------------


@mcp.tool(annotations=MUTATING)
def ppt_set_style_profile(
    name: str,
    system_prompt: Annotated[str, Field(description="The corporate visual-language prompt")],
    metadata: Annotated[
        dict[str, Any] | None,
        Field(description="Optional: allowed_colors, forbidden_motifs, composition, "
                          "media_type, negative_prompt_base, default_aspect_ratio, "
                          "text_in_image, logo_usage"),
    ] = None,
) -> dict[str, Any]:
    """Create or update a brand style profile (local JSON, versioned on every
    update). Profiles govern image-placeholder prompts; they are configuration,
    never executable instruction."""
    profile = styles.set(name, system_prompt, metadata)
    return {"profile_id": profile["profile_id"], "version": profile["version"]}


@mcp.tool(annotations=READ_ONLY)
def ppt_get_style_profile(profile_id: str) -> dict[str, Any]:
    """Full style profile including version history."""
    return styles.get(profile_id)


@mcp.tool(annotations=READ_ONLY)
def ppt_list_style_profiles() -> dict[str, Any]:
    """List brand style profiles."""
    return {
        "profiles": [
            {"profile_id": p["profile_id"], "name": p["name"], "version": p["version"]}
            for p in styles.list()
        ]
    }


@mcp.tool(annotations=DESTRUCTIVE)
def ppt_delete_style_profile(profile_id: str) -> dict[str, Any]:
    """Delete a style profile. Prompts already embedded in decks keep their
    recorded profile id/version for provenance."""
    styles.delete(profile_id)
    return {"profile_id": profile_id, "deleted": True}


def _slide_context(prs, slide_index: int) -> dict[str, Any]:
    detail = reader.slide_detail(prs, slide_index)
    bullets = [
        p["text"]
        for s in detail["shapes"]
        if s["placeholder"] and s["placeholder"]["role"] == "body" and s["text"]
        for p in s["text"]
        if p["text"].strip()
    ]
    return {
        "title": detail["title"],
        "body_summary": "; ".join(bullets[:4]) if bullets else None,
        "layout": detail["layout"],
    }


@mcp.tool(annotations=READ_ONLY)
def ppt_compose_image_prompt(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    image_intent: str,
    profile_id: str,
    constraints: str | None = None,
    aspect_ratio: str | None = None,
) -> dict[str, Any]:
    """Compose a governed image-generation prompt from the slide's content and
    a brand style profile — deterministic assembly; refine the scene wording
    if needed, then store it with ppt_create_image_placeholder."""
    profile = styles.get(profile_id)
    bundle = StyleProfileRegistry.compose_prompt(
        profile, image_intent, _slide_context(_prs(deck_id), slide_index),
        aspect_ratio=aspect_ratio, constraints=constraints,
    )
    return {"deck_id": deck_id, "slide_index": slide_index, **bundle}


@mcp.tool(annotations=MUTATING)
def ppt_create_image_placeholder(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    spec: ImagePlaceholderSpec,
    dry_run: bool = False,
) -> dict[str, Any]:
    """Create a governed image slot: target a picture placeholder (preferred)
    or, with allow_freeform=true + position, a labeled box. The full prompt is
    stored in the in-file manifest + a pointer in the speaker notes — never on
    the visible slide."""
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)

    profile = styles.get(spec.profile_id) if spec.profile_id else None
    if spec.prompt is not None:
        bundle = {
            "prompt": spec.prompt,
            "negative_prompt": spec.negative_prompt
            or (profile or {}).get("negative_prompt_base", ""),
            "alt_text": spec.alt_text or spec.image_intent,
            "aspect_ratio": spec.aspect_ratio
            or (profile or {}).get("default_aspect_ratio", "16:9"),
            "provenance": {
                "style_profile_id": profile["profile_id"] if profile else None,
                "style_profile_version": profile["version"] if profile else None,
                "created_by_tool": "ppt_create_image_placeholder",
            },
        }
    elif profile is not None:
        bundle = StyleProfileRegistry.compose_prompt(
            profile, spec.image_intent, _slide_context(prs, slide_index),
            aspect_ratio=spec.aspect_ratio,
        )
        if spec.alt_text:
            bundle["alt_text"] = spec.alt_text
    else:
        raise PptMcpError(
            "Provide either a prompt or a profile_id (ppt_list_style_profiles) "
            "so the prompt can be composed."
        )
    if profile is not None:
        violations = StyleProfileRegistry.validate_prompt(profile, bundle["prompt"])
        if violations:
            raise PptMcpError(
                "Prompt violates the style profile: " + "; ".join(violations)
                + ". Adjust the prompt or the profile."
            )

    if spec.shape_ref is not None:
        target = writer.resolve_shape(slide, spec.shape_ref)
        if not hasattr(target, "insert_picture"):
            raise PptMcpError(
                f"Shape '{target.name}' is not a picture placeholder. Target a "
                "placeholder with role 'picture' (ppt_get_slide lists roles), or "
                "use allow_freeform=true with a position."
            )
        target_kind = "picture_placeholder"
        plan = {"target": f"picture placeholder shape_id {target.shape_id}"}
    elif spec.allow_freeform and spec.position is not None:
        target = None
        target_kind = "freeform_label"
        plan = {"target": "freeform labeled box", "warnings": [writer.FREEFORM_WARNING]}
    else:
        raise PptMcpError(
            "Target a picture placeholder via shape_ref, or pass "
            "allow_freeform=true with a position."
        )
    plan.update(
        {"slide_index": slide_index, "image_intent": spec.image_intent,
         "prompt_chars": len(bundle["prompt"])}
    )
    if dry_run:
        return {"applied": False, "plan": plan, "prompt_bundle": bundle}

    if target is None:
        from pptx.util import Inches

        pos = spec.position
        target = slide.shapes.add_textbox(
            Inches(pos.left), Inches(pos.top), Inches(pos.width), Inches(pos.height)
        )
        target.text_frame.text = f"Image placeholder: {spec.image_intent}"

    record = {
        "slide_index": slide_index,
        "shape_id": target.shape_id,
        "target_kind": target_kind,
        "image_intent": spec.image_intent,
        "prompt": bundle["prompt"],
        "negative_prompt": bundle["negative_prompt"],
        "aspect_ratio": bundle["aspect_ratio"],
        "alt_text": bundle["alt_text"],
        "status": "pending",
        "layout": slide.slide_layout.name,
        **bundle["provenance"],
    }
    data = deck_manifest.load(prs)
    data["image_placeholders"].append(record)
    deck_manifest.save(prs, data)

    pointer = f"[ppt-mcp] Image placeholder (shape {target.shape_id}): {spec.image_intent} — full prompt in deck manifest."
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = (notes_frame.text + "\n" + pointer).strip()

    _commit(deck_id, prs, "ppt_create_image_placeholder",
            slide_index=slide_index, shape_id=target.shape_id)
    return {"applied": True, "plan": plan, "shape_id": target.shape_id,
            "record": record}


@mcp.tool(annotations=READ_ONLY)
def ppt_list_image_placeholders(
    deck_id: str,
    status: Annotated[
        str | None, Field(description="Filter: pending | generated | approved")
    ] = None,
) -> dict[str, Any]:
    """All governed image slots in the deck, from the in-file manifest."""
    prs = _prs(deck_id)
    data = deck_manifest.load(prs)
    deck_manifest.relocate_records(data, prs)
    records = data["image_placeholders"]
    if status is not None:
        records = [r for r in records if r["status"] == status]
    return {"deck_id": deck_id, "count": len(records), "image_placeholders": records}


@mcp.tool(annotations=MUTATING)
def ppt_update_image_placeholder(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    shape_ref: str,
    patch: Annotated[
        dict[str, Any],
        Field(description="Editable: prompt, negative_prompt, alt_text, image_intent, "
                          "aspect_ratio, status"),
    ],
    dry_run: bool = False,
) -> dict[str, Any]:
    """Update a governed image slot's manifest record (e.g. a refined prompt).
    The prompt is re-validated against the recorded style profile."""
    allowed = {"prompt", "negative_prompt", "alt_text", "image_intent", "aspect_ratio", "status"}
    unknown = set(patch) - allowed
    if unknown:
        raise PptMcpError(f"Cannot patch {sorted(unknown)}. Editable: {sorted(allowed)}.")
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)
    shape = writer.resolve_shape(slide, shape_ref)
    data = deck_manifest.load(prs)
    record = deck_manifest.find_record(data, slide_index, shape.shape_id)
    if record is None:
        raise PptMcpError(
            f"No image-placeholder record for shape {shape.shape_id} on slide "
            f"{slide_index}. ppt_list_image_placeholders shows what exists."
        )
    if "prompt" in patch and record.get("style_profile_id"):
        profile = styles.get(record["style_profile_id"])
        violations = StyleProfileRegistry.validate_prompt(profile, patch["prompt"])
        if violations:
            raise PptMcpError(
                "Updated prompt violates the style profile: " + "; ".join(violations)
            )
    if dry_run:
        return {"applied": False, "plan": {"update": sorted(patch)}, "current": record}
    record.update(patch)
    deck_manifest.save(prs, data)
    _commit(deck_id, prs, "ppt_update_image_placeholder",
            slide_index=slide_index, shape_id=shape.shape_id)
    return {"applied": True, "record": record}


@mcp.tool(annotations=MUTATING)
def ppt_fill_image_placeholder(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    shape_ref: str,
    image_path: Annotated[str, Field(description="The externally generated image file")],
    dry_run: bool = False,
) -> dict[str, Any]:
    """Insert the generated image into a governed slot: picture placeholders
    get insert_picture; freeform label boxes are replaced in place. The
    manifest record flips to status 'generated'."""
    prs = _prs(deck_id)
    slide = reader.get_slide_by_index(prs, slide_index)
    shape = writer.resolve_shape(slide, shape_ref)
    image_file = Path(image_path).expanduser()
    if not image_file.is_file():
        raise PptMcpError(f"Image file not found: '{image_file}'.")
    data = deck_manifest.load(prs)
    record = deck_manifest.find_record(data, slide_index, shape.shape_id)
    if dry_run:
        return {"applied": False,
                "plan": {"fill": shape.shape_id, "image": str(image_file),
                         "has_manifest_record": record is not None}}
    if hasattr(shape, "insert_picture"):
        picture = shape.insert_picture(str(image_file))
    else:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        shape._element.getparent().remove(shape._element)
        picture = slide.shapes.add_picture(
            str(image_file), left, top, width=width, height=height
        )
    if record is not None:
        if record.get("alt_text"):
            picture._element._nvXxPr.cNvPr.set("descr", record["alt_text"])
        record.update(
            {"status": "generated", "shape_id": picture.shape_id,
             "image_path": str(image_file)}
        )
        deck_manifest.save(prs, data)
    _commit(deck_id, prs, "ppt_fill_image_placeholder",
            slide_index=slide_index, shape_id=picture.shape_id)
    return {"applied": True, "shape_id": picture.shape_id,
            "status": "generated" if record else "filled (no manifest record)"}


# -- compliance (§10.5) ------------------------------------------------------------


def _template_entry_or_none(template_id: str | None) -> dict[str, Any] | None:
    return registry.get(template_id) if template_id is not None else None


@mcp.tool(annotations=READ_ONLY)
def ppt_validate_compliance(
    deck_id: str,
    template_id: Annotated[
        str | None,
        Field(description="Reference template. Omit to validate against the deck's own theme "
                          "(template-relative rules C01/C08/C10 are then skipped)."),
    ] = None,
) -> dict[str, Any]:
    """Check the deck against compliance rules C01–C10: layout provenance,
    theme fonts/colors, footers, covered logos, placeholder bypass, estimated
    text overflow, density, off-grid shapes, slide size. Returns structured
    findings with severity and auto_fixable flags."""
    findings = compliance.validate(_prs(deck_id), _template_entry_or_none(template_id))
    summary = {
        level: len([f for f in findings if f["severity"] == level])
        for level in ("error", "warning", "info")
    }
    return {"deck_id": deck_id, "template_id": template_id, "summary": summary,
            "findings": findings}


@mcp.tool(annotations=MUTATING)
def ppt_repair_compliance(
    deck_id: str,
    template_id: str | None = None,
    strategy: Annotated[str, Field(description="Only 'conservative' for now")] = "conservative",
    dry_run: Annotated[bool, Field(description="Default TRUE — repairs must be previewed")] = True,
) -> dict[str, Any]:
    """Fix auto-fixable compliance findings. Conservative strategy: re-link
    non-theme fonts to the theme, snap near-theme colors to exact theme values.
    dry_run defaults to true; pass dry_run=false to apply."""
    prs = _prs(deck_id)
    entry = _template_entry_or_none(template_id)
    findings = compliance.validate(prs, entry)
    fixes = compliance.plan_repairs(prs, findings, entry, strategy=strategy)
    if dry_run:
        return {"applied": False, "planned_fixes": fixes,
                "unfixable_findings": len(findings) - len(fixes)}
    applied = compliance.apply_repairs(prs, fixes, entry)
    _commit(deck_id, prs, "ppt_repair_compliance", fixes=len(fixes))
    remaining = compliance.validate(_prs(deck_id), entry)
    return {
        "applied": True,
        "fixes_applied": applied,
        "remaining_findings": len(remaining),
        "remaining_by_severity": {
            level: len([f for f in remaining if f["severity"] == level])
            for level in ("error", "warning", "info")
        },
    }


@mcp.tool(annotations=DESTRUCTIVE)
def ppt_apply_template(
    deck_id: str,
    template_id: str,
    dry_run: Annotated[
        bool,
        Field(description="Default TRUE — always review the fidelity report before applying"),
    ] = True,
) -> dict[str, Any]:
    """Re-target a deck onto a registered template: each slide is re-created
    from the best-matching layout with its placeholder content migrated;
    freeform/chart shapes are carried over unchanged and flagged. The dry-run
    fidelity report enumerates per-slide risks. Applying snapshots first
    (ppt_undo recovers) and runs compliance validation automatically."""
    prs = _prs(deck_id)
    entry = registry.get(template_id)
    session = sessions.get(deck_id)
    assets_dir = session.session_dir / "migration_assets"
    new_prs, plan = retarget.retarget_deck(
        prs, entry, registry.materialized_path(template_id), assets_dir, apply=not dry_run
    )
    if dry_run:
        return {"applied": False, "template_id": template_id, "plan": plan}
    sessions.snapshot(deck_id)
    new_prs.save(str(session.working_path))
    store.log_provenance("ppt_apply_template", deck_id=deck_id, template_id=template_id)
    findings = compliance.validate(_prs(deck_id), entry)
    return {
        "applied": True,
        "template_id": template_id,
        "plan": plan,
        "validation_summary": {
            level: len([f for f in findings if f["severity"] == level])
            for level in ("error", "warning", "info")
        },
        "validation_findings": findings,
    }


@mcp.tool(annotations=MUTATING)
def ppt_extract_template_from_deck(
    deck_id: str,
    name: Annotated[str | None, Field(description="Name for the derived template")] = None,
    version: str | None = None,
) -> dict[str, Any]:
    """Register an open deck's design system as a *derived* template (for the
    'Final_v13_really_final.pptx' case where no clean .potx exists). Captures
    masters, layouts, theme, and per-layout usage statistics from the deck's
    slides. Curate the inferred intent tags afterwards with ppt_update_template."""
    session = sessions.get(deck_id)
    entry = registry.register(
        session.working_path,
        name=name or f"Derived from {deck_id}",
        version=version,
        derived=True,
    )
    return {
        "already_registered": entry.get("already_registered", False),
        "template_id": entry["template_id"],
        "name": entry["name"],
        "derived": True,
        "layouts": len(entry["layouts"]),
        "layout_usage": entry["layout_usage"],
    }


# -- rendering & export (§10.7) ------------------------------------------------------


@mcp.tool(annotations=READ_ONLY)
def ppt_render_slide(
    deck_id: str,
    slide_index: Annotated[int, Field(ge=1)],
    dpi: Annotated[int, Field(ge=36, le=300)] = 96,
):
    """Render one slide to PNG via headless LibreOffice and return the image.
    Treat renders as validation evidence — PowerPoint is the fidelity arbiter."""
    from mcp.server.fastmcp import Image

    session = sessions.get(deck_id)
    rendered = renderer.render_slides(session.working_path, [slide_index], dpi=dpi)
    png_path = rendered[slide_index]
    return [Image(path=str(png_path)), f"Rendered slide {slide_index} -> {png_path}"]


@mcp.tool(annotations=READ_ONLY)
def ppt_render_deck(
    deck_id: str,
    dpi: Annotated[int, Field(ge=36, le=300)] = 96,
) -> dict[str, Any]:
    """Render every slide to PNG; returns the file paths (use ppt_render_slide
    to view one inline)."""
    session = sessions.get(deck_id)
    rendered = renderer.render_slides(session.working_path, None, dpi=dpi)
    return {"deck_id": deck_id, "dpi": dpi,
            "slides": {str(i): str(p) for i, p in sorted(rendered.items())}}


@mcp.tool(annotations=MUTATING)
def ppt_export_pdf(
    deck_id: str,
    path: Annotated[str, Field(description="Target .pdf path")],
) -> dict[str, Any]:
    """Export the deck to PDF (headless LibreOffice)."""
    import shutil as _shutil

    target = Path(path).expanduser().resolve()
    if target.suffix.lower() != ".pdf":
        raise PptMcpError(f"Target must end in .pdf, got '{target.name}'.")
    session = sessions.get(deck_id)
    pdf = renderer.to_pdf(session.working_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    _shutil.copy2(pdf, target)
    store.log_provenance("ppt_export_pdf", deck_id=deck_id, target=target)
    return {"deck_id": deck_id, "exported_to": str(target)}


@mcp.tool(annotations=READ_ONLY)
def ppt_visual_diff(
    deck_id: str,
    snapshot: Annotated[
        int, Field(ge=1, description="Snapshot number (1 = oldest); see undo_steps_available"),
    ],
    dpi: Annotated[int, Field(ge=36, le=300)] = 96,
) -> dict[str, Any]:
    """Pixel-diff the current deck against one of its pre-mutation snapshots —
    verifies e.g. that a logo or footer survived an edit."""
    session = sessions.get(deck_id)
    snapshot_path = session.snapshots_dir / f"{snapshot}.pptx"
    if not snapshot_path.is_file():
        raise PptMcpError(
            f"Snapshot {snapshot} does not exist for deck '{deck_id}' "
            f"({session.snapshot_count} snapshot(s) available)."
        )
    current = renderer.render_slides(session.working_path, None, dpi=dpi)
    previous = renderer.render_slides(snapshot_path, None, dpi=dpi)
    slides: dict[str, Any] = {}
    for index in sorted(set(current) | set(previous)):
        if index not in current:
            slides[str(index)] = {"changed": True, "note": "slide removed"}
        elif index not in previous:
            slides[str(index)] = {"changed": True, "note": "slide added"}
        else:
            slides[str(index)] = diff_images(previous[index], current[index])
    changed = [i for i, d in slides.items() if d.get("changed")]
    return {"deck_id": deck_id, "snapshot": snapshot,
            "changed_slides": changed, "slides": slides}


def main() -> None:
    log.info("ppt_mcp starting (stdio); data dir: %s", store.root)
    mcp.run()


if __name__ == "__main__":
    main()

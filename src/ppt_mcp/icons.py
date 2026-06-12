"""Icon library (DESIGN.md §9.3, scoped): vendored Material Symbols set +
icon sets harvested from existing PowerPoint files. Fully offline at runtime.

Rendering: SVG icons are rasterized black-on-white (svglib/reportlab), then
the inverted grayscale becomes the alpha channel of a solid-color image —
clean anti-aliased transparent PNGs tinted in any theme color, no SVG
surgery. Raster icons (typically harvested) are inserted as-is.
"""

from __future__ import annotations

import hashlib
import io
import json
import re
from pathlib import Path
from typing import Any

from PIL import Image, ImageOps
from rapidfuzz import fuzz

from .errors import PptMcpError
from .store import Store

BUILTIN_DIR = Path(__file__).parent / "assets" / "icons"

RASTER_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp"}
HARVEST_EXTS = RASTER_EXTS | {".svg"}


def _slug(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_") or "icon"


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    if not re.fullmatch(r"[0-9a-fA-F]{6}", value):
        raise PptMcpError(f"'{value}' is not a #RRGGBB color.")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def _render_svg_tinted(svg_path: Path, color_hex: str, px: int) -> Image.Image:
    from reportlab.graphics import renderPM
    from svglib.svglib import svg2rlg

    drawing = svg2rlg(str(svg_path))
    if drawing is None:
        raise PptMcpError(f"Could not parse SVG '{svg_path.name}'.")
    scale = px / max(drawing.width, drawing.height)
    drawing.width *= scale
    drawing.height *= scale
    drawing.scale(scale, scale)
    rendered = renderPM.drawToPIL(drawing, bg=0xFFFFFF)
    alpha = ImageOps.invert(rendered.convert("L"))
    tinted = Image.new("RGBA", rendered.size, (*_hex_to_rgb(color_hex), 255))
    tinted.putalpha(alpha)
    return tinted


class IconLibrary:
    def __init__(self, store: Store):
        self.store = store
        self.user_dir = store.assets_dir / "icons"
        self.cache_dir = store.renders_dir / "icons"

    # -- sets -----------------------------------------------------------------

    def _set_dirs(self) -> dict[str, Path]:
        dirs: dict[str, Path] = {}
        for base in (BUILTIN_DIR, self.user_dir):
            if not base.is_dir():
                continue
            for index_path in base.glob("*/index.json"):
                set_dir = index_path.parent
                dirs[set_dir.name] = set_dir  # user sets shadow builtins by id
        return dirs

    def _load_index(self, set_dir: Path) -> dict[str, Any]:
        return json.loads((set_dir / "index.json").read_text(encoding="utf-8"))

    def list_sets(self) -> list[dict[str, Any]]:
        sets = []
        for set_id, set_dir in sorted(self._set_dirs().items()):
            index = self._load_index(set_dir)
            sets.append(
                {
                    "set_id": set_id,
                    "name": index.get("name", set_id),
                    "license": index.get("license", "unknown"),
                    "icon_count": len(index["icons"]),
                    "recolorable": index.get("recolorable", False),
                    "builtin": set_dir.parent == BUILTIN_DIR,
                }
            )
        return sets

    def _get_set(self, set_id: str) -> tuple[Path, dict[str, Any]]:
        dirs = self._set_dirs()
        if set_id not in dirs:
            known = ", ".join(sorted(dirs)) or "none"
            raise PptMcpError(
                f"Unknown icon set '{set_id}'. Available sets: {known}. Harvest one "
                "from a PowerPoint file with ppt_harvest_icons."
            )
        return dirs[set_id], self._load_index(dirs[set_id])

    # -- search & resolve --------------------------------------------------------

    def search(
        self, query: str, set_id: str | None = None, limit: int = 12
    ) -> list[dict[str, Any]]:
        targets = [set_id] if set_id else list(self._set_dirs())
        results = []
        for sid in targets:
            _, index = self._get_set(sid)
            for icon in index["icons"]:
                score = fuzz.WRatio(query.lower(), icon["name"].lower())
                for tag in icon.get("tags", []):
                    score = max(score, fuzz.WRatio(query.lower(), tag.lower()))
                if score >= 55:
                    results.append(
                        {
                            "set_id": sid,
                            "icon_id": icon["id"],
                            "name": icon["name"],
                            "tags": icon.get("tags", []),
                            "format": icon["format"],
                            "score": round(score / 100, 2),
                        }
                    )
        results.sort(key=lambda r: r["score"], reverse=True)
        return results[:limit]

    def get_icon(self, set_id: str, icon_id: str) -> tuple[dict[str, Any], Path]:
        set_dir, index = self._get_set(set_id)
        for icon in index["icons"]:
            if icon["id"] == icon_id:
                return icon, set_dir / icon["file"]
        raise PptMcpError(
            f"No icon '{icon_id}' in set '{set_id}'. Use ppt_search_icons to find one."
        )

    # -- rendering ----------------------------------------------------------------

    def render_png(
        self, set_id: str, icon_id: str, color_hex: str, px: int
    ) -> tuple[Path, bool]:
        """Render an icon to a PNG file. Returns (path, was_tinted)."""
        icon, source = self.get_icon(set_id, icon_id)
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        if icon["format"] == "svg":
            safe_color = color_hex.lstrip("#").upper()
            out = self.cache_dir / f"{set_id}_{icon_id}_{safe_color}_{px}.png"
            if not out.is_file():
                _render_svg_tinted(source, color_hex, px).save(out)
            return out, True
        return source, False  # raster icons go in as-is; tinting would mangle them

    # -- harvesting ------------------------------------------------------------------

    def harvest_from_pptx(
        self,
        pptx_path: Path,
        set_name: str,
        min_px: int = 16,
        max_px: int = 600,
    ) -> dict[str, Any]:
        """Import icon-sized images from a PowerPoint file into a local set.
        Raster images outside [min_px, max_px] (longest edge) are skipped;
        SVG media are always taken. Duplicates (by content hash) are skipped."""
        from pptx import Presentation

        prs = Presentation(str(pptx_path))

        # image content-hash -> shape name(s), for human-readable icon names
        names_by_sha: dict[str, str] = {}
        for slide in prs.slides:
            for shape in slide.shapes:
                image = getattr(shape, "image", None)
                try:
                    sha = image.sha1 if image is not None else None
                except (AttributeError, ValueError):
                    sha = None
                if sha and sha not in names_by_sha:
                    names_by_sha[sha] = shape.name

        set_id = _slug(set_name)
        set_dir = self.user_dir / set_id
        icons_dir = set_dir / "icons"
        icons_dir.mkdir(parents=True, exist_ok=True)
        if (set_dir / "index.json").is_file():
            index = self._load_index(set_dir)
        else:
            index = {
                "set_id": set_id,
                "name": set_name,
                "license": f"harvested from {pptx_path.name} — clear rights before reuse",
                "recolorable": False,
                "icons": [],
            }
        known_hashes = {icon.get("sha1") for icon in index["icons"]}
        taken_ids = {icon["id"] for icon in index["icons"]}

        imported, skipped_size, skipped_dupe = [], 0, 0
        for part in prs.part.package.iter_parts():
            partname = str(part.partname)
            ext = Path(partname).suffix.lower()
            if not partname.startswith("/ppt/media/") or ext not in HARVEST_EXTS:
                continue
            blob = part.blob
            sha = hashlib.sha1(blob).hexdigest()
            if sha in known_hashes:
                skipped_dupe += 1
                continue
            if ext in RASTER_EXTS:
                try:
                    with Image.open(io.BytesIO(blob)) as im:
                        longest = max(im.size)
                except OSError:
                    continue
                if not min_px <= longest <= max_px:
                    skipped_size += 1
                    continue
            base = _slug(names_by_sha.get(sha, Path(partname).stem))
            icon_id = base
            counter = 2
            while icon_id in taken_ids:
                icon_id = f"{base}_{counter}"
                counter += 1
            taken_ids.add(icon_id)
            known_hashes.add(sha)
            file_name = f"icons/{icon_id}{ext}"
            (set_dir / file_name).write_bytes(blob)
            entry = {
                "id": icon_id,
                "name": icon_id.replace("_", " "),
                "tags": [],
                "file": file_name,
                "format": "svg" if ext == ".svg" else ext.lstrip("."),
                "sha1": sha,
                "source": pptx_path.name,
            }
            index["icons"].append(entry)
            imported.append(icon_id)

        (set_dir / "index.json").write_text(
            json.dumps(index, indent=1, ensure_ascii=False), encoding="utf-8"
        )
        self.store.log_provenance(
            "ppt_harvest_icons", set_id=set_id, source=pptx_path, imported=len(imported)
        )
        return {
            "set_id": set_id,
            "imported": imported,
            "imported_count": len(imported),
            "skipped_wrong_size": skipped_size,
            "skipped_duplicates": skipped_dupe,
            "total_in_set": len(index["icons"]),
        }

"""Read-side OOXML engine (DESIGN.md §10.2): deck overviews, per-slide shape
inventories with stable addressing (§4.2), and text search.

Slide indices are 1-based at the API boundary — matching the PowerPoint UI —
and converted at exactly one place (`get_slide_by_index`).
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from .errors import SlideIndexError

EMU_PER_INCH = 914400

# Placeholder type name -> semantic role (DESIGN.md §4.2).
ROLE_BY_PH_TYPE_NAME = {
    "TITLE": "title",
    "CENTER_TITLE": "title",
    "VERTICAL_TITLE": "title",
    "SUBTITLE": "subtitle",
    "BODY": "body",
    "VERTICAL_BODY": "body",
    "OBJECT": "body",
    "VERTICAL_OBJECT": "body",
    "PICTURE": "picture",
    "BITMAP": "picture",
    "CLIP_ART": "picture",
    "CHART": "chart",
    "TABLE": "table",
    "FOOTER": "footer",
    "SLIDE_NUMBER": "slide_number",
    "DATE": "date",
    "HEADER": "header",
    "MEDIA_CLIP": "media",
    "ORG_CHART": "diagram",
}


def load_presentation(path: Path) -> PresentationType:
    return Presentation(str(path))


def emu_to_inches(emu: int | None) -> float | None:
    return round(emu / EMU_PER_INCH, 3) if emu is not None else None


def get_slide_by_index(prs: PresentationType, slide_index: int) -> Slide:
    """`slide_index` is 1-based."""
    count = len(prs.slides)
    if not 1 <= slide_index <= count:
        raise SlideIndexError(slide_index, count)
    return prs.slides[slide_index - 1]


def _aspect_ratio(width_emu: int, height_emu: int) -> str:
    for w, h in ((16, 9), (4, 3), (16, 10)):
        if width_emu * h == height_emu * w:
            return f"{w}:{h}"
    return f"custom ({round(width_emu / height_emu, 3)})"


def _placeholder_role(ph_format) -> str:
    type_name = ph_format.type.name if ph_format.type is not None else "BODY"
    return ROLE_BY_PH_TYPE_NAME.get(type_name, type_name.lower())


def _shape_text_paragraphs(shape) -> list[dict[str, Any]] | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    return [
        {"text": p.text, "level": p.level}
        for p in shape.text_frame.paragraphs
    ]


def _inherited_geometry(shape, slide: Slide) -> dict | None:
    """Placeholders without explicit position inherit it from the layout (or
    master). Resolve one level up so the agent always sees real geometry."""
    ph_idx = shape.placeholder_format.idx
    for scope in (slide.slide_layout, slide.slide_layout.slide_master):
        for layout_ph in scope.placeholders:
            if layout_ph.placeholder_format.idx == ph_idx and layout_ph.left is not None:
                return _geometry_dict(layout_ph, inherited=True)
    return None


def _geometry_dict(shape, inherited: bool = False) -> dict | None:
    if shape.left is None:
        return None
    geometry = {
        "left_in": emu_to_inches(shape.left),
        "top_in": emu_to_inches(shape.top),
        "width_in": emu_to_inches(shape.width),
        "height_in": emu_to_inches(shape.height),
        "left_emu": int(shape.left),
        "top_emu": int(shape.top),
        "width_emu": int(shape.width) if shape.width is not None else None,
        "height_emu": int(shape.height) if shape.height is not None else None,
    }
    if inherited:
        geometry["inherited_from_layout"] = True
    return geometry


def shape_info(shape, slide: Slide) -> dict[str, Any]:
    """One shape, with all three address forms echoed (DESIGN.md §4.2)."""
    info: dict[str, Any] = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": shape.shape_type.name if shape.shape_type is not None else None,
        "placeholder": None,
        "geometry": _geometry_dict(shape),
        "text": _shape_text_paragraphs(shape),
    }
    if shape.is_placeholder:
        ph = shape.placeholder_format
        info["placeholder"] = {
            "idx": ph.idx,
            "type": ph.type.name if ph.type is not None else None,
            "role": _placeholder_role(ph),
        }
        if info["geometry"] is None:
            info["geometry"] = _inherited_geometry(shape, slide)
    if getattr(shape, "has_table", False):
        table = shape.table
        info["table"] = {"rows": len(table.rows), "columns": len(table.columns)}
    if getattr(shape, "has_chart", False):
        info["chart"] = {"chart_type": shape.chart.chart_type.name}
    return info


def slide_title(slide: Slide) -> str | None:
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()
        return text or None
    return None


def slide_notes(slide: Slide) -> str | None:
    if slide.has_notes_slide:
        text = slide.notes_slide.notes_text_frame.text.strip()
        return text or None
    return None


def deck_overview(prs: PresentationType) -> dict[str, Any]:
    width, height = int(prs.slide_width), int(prs.slide_height)
    return {
        "slide_count": len(prs.slides),
        "slide_size": {
            "width_in": emu_to_inches(width),
            "height_in": emu_to_inches(height),
            "aspect_ratio": _aspect_ratio(width, height),
        },
        "masters": [
            {
                "name": master.name or f"Master {i + 1}",
                "layouts": [layout.name for layout in master.slide_layouts],
            }
            for i, master in enumerate(prs.slide_masters)
        ],
        "slides": [
            {
                "slide_index": i + 1,
                "layout": slide.slide_layout.name,
                "title": slide_title(slide),
                "has_notes": slide.has_notes_slide and bool(slide_notes(slide)),
            }
            for i, slide in enumerate(prs.slides)
        ],
    }


def slide_detail(prs: PresentationType, slide_index: int) -> dict[str, Any]:
    slide = get_slide_by_index(prs, slide_index)
    return {
        "slide_index": slide_index,
        "layout": slide.slide_layout.name,
        "master": slide.slide_layout.slide_master.name or None,
        "title": slide_title(slide),
        "shapes": [shape_info(shape, slide) for shape in slide.shapes],
        "notes": slide_notes(slide),
    }


def search_deck(
    prs: PresentationType, query: str, limit: int = 25, offset: int = 0
) -> dict[str, Any]:
    """Case-insensitive substring search over shape text and notes."""
    needle = query.lower()
    hits: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if needle in paragraph.text.lower():
                    hits.append(
                        {
                            "slide_index": i + 1,
                            "shape_id": shape.shape_id,
                            "shape_name": shape.name,
                            "match": paragraph.text,
                            "in": "shape",
                        }
                    )
        notes = slide_notes(slide)
        if notes and needle in notes.lower():
            line = next(
                (ln for ln in notes.splitlines() if needle in ln.lower()), notes
            )
            hits.append(
                {
                    "slide_index": i + 1,
                    "shape_id": None,
                    "shape_name": "notes",
                    "match": line,
                    "in": "notes",
                }
            )
    page = hits[offset : offset + limit]
    return {
        "query": query,
        "total_count": len(hits),
        "count": len(page),
        "offset": offset,
        "has_more": offset + limit < len(hits),
        "next_offset": offset + limit if offset + limit < len(hits) else None,
        "hits": page,
    }

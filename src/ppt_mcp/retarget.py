"""Template application engine (DESIGN.md §7.4) — the riskiest operation in
the product, so it is plan-first by design.

Re-targeting walks every source slide, extracts its semantic content
(title/subtitle/body/table/picture from placeholders), chooses the best
target layout (exact name match, else the §5.3 scorer), and re-creates the
slide on the new template's design system. Everything that cannot be
expressed as placeholder content — freeform shapes, charts, extra pictures —
is carried over byte-faithfully with relationship remapping and *flagged*,
never dropped silently.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide

from . import recommend
from .models import ContentSpec, ImageRef, Paragraph, TableSpec
from .reader import emu_to_inches, load_presentation, slide_notes
from .writer import (
    _SHAPE_TAGS,
    apply_add_slide,
    clone_shape_elements,
    delete_slide,
    placeholder_role,
    plan_content_mapping,
    resolve_layout,
)


def _extract_slide(
    slide: Slide, slide_index: int, assets_dir: Path | None
) -> tuple[ContentSpec, list, list[str], list[str]]:
    """Pull semantic content out of a slide's placeholders.

    Returns (spec, orphan_elements, carried_over_names, warnings). With
    assets_dir=None (planning), picture blobs are not written to disk."""
    title = subtitle = None
    body: list[Paragraph] = []
    table: TableSpec | None = None
    image: ImageRef | None = None
    warnings: list[str] = []
    handled_elements: list = []

    for ph in slide.placeholders:
        role = placeholder_role(ph)
        if role in ("footer", "date", "slide_number"):
            handled_elements.append(ph._element)  # the new template's own rules apply
            continue
        if role == "title" and getattr(ph, "has_text_frame", False):
            title = ph.text_frame.text.strip() or None
            handled_elements.append(ph._element)
        elif role == "subtitle" and getattr(ph, "has_text_frame", False):
            subtitle = ph.text_frame.text.strip() or None
            handled_elements.append(ph._element)
        elif getattr(ph, "has_table", False):
            if table is None:
                rows = [
                    [cell.text for cell in row.cells] for row in ph.table.rows
                ]
                table = TableSpec(rows=rows) if rows else None
                handled_elements.append(ph._element)
        elif role == "picture":
            picture_image = getattr(ph, "image", None)
            try:
                blob, ext = (picture_image.blob, picture_image.ext) if picture_image else (None, None)
            except (AttributeError, ValueError):
                blob = None
            if blob is None:
                handled_elements.append(ph._element)  # empty picture placeholder
            elif image is None:
                if assets_dir is not None:
                    assets_dir.mkdir(parents=True, exist_ok=True)
                    target = assets_dir / f"slide{slide_index}-{ph.shape_id}.{ext}"
                    target.write_bytes(blob)
                    image = ImageRef(path=str(target))
                else:
                    image = ImageRef(path="(extracted at apply)")
                handled_elements.append(ph._element)
        elif role == "body" and getattr(ph, "has_text_frame", False):
            if not ph.text_frame.text.strip():
                handled_elements.append(ph._element)
                continue
            paragraphs = [
                Paragraph(text=p.text, level=p.level)
                for p in ph.text_frame.paragraphs
                if p.text.strip()
            ]
            if body:
                warnings.append(
                    "multiple filled body placeholders merged into one body section"
                )
            body.extend(paragraphs)
            handled_elements.append(ph._element)

    handled_ids = {id(el) for el in handled_elements}
    orphans = [
        child
        for child in slide.shapes._spTree
        if child.tag in _SHAPE_TAGS and id(child) not in handled_ids
    ]
    orphan_ids = {id(o) for o in orphans}
    orphan_names = []
    for shape in slide.shapes:
        if id(shape._element) not in orphan_ids:
            continue
        label = shape.name
        text = (
            shape.text_frame.text.strip()
            if getattr(shape, "has_text_frame", False)
            else ""
        )
        if text:
            snippet = text if len(text) <= 40 else text[:37] + "..."
            label += f" ('{snippet}')"
        orphan_names.append(label)
    spec = ContentSpec(
        title=title, subtitle=subtitle, body=body or None, table=table,
        image=image, notes=slide_notes(slide),
    )
    return spec, orphans, orphan_names, warnings


def _choose_layout(
    template_entry: dict[str, Any], source_layout_name: str, spec: ContentSpec
) -> tuple[dict[str, Any], str]:
    for layout in template_entry["layouts"]:
        if layout["name"].lower() == source_layout_name.lower():
            return layout, "exact"
    best = recommend.recommend_layouts(
        template_entry["layouts"], source_layout_name, spec, top_n=1
    )[0]
    target = next(
        lo for lo in template_entry["layouts"] if lo["layout_id"] == best["layout_id"]
    )
    return target, f"scored {best['confidence']} ({best['reason']})"


def retarget_deck(
    prs_source: PresentationType,
    template_entry: dict[str, Any],
    materialized_path: Path,
    assets_dir: Path | None,
    apply: bool,
) -> tuple[PresentationType | None, dict[str, Any]]:
    """Plan (apply=False) or perform (apply=True) the re-targeting.
    Returns (new_presentation_or_None, fidelity_plan)."""
    new_prs = load_presentation(materialized_path) if apply else None
    if new_prs is not None:
        for index in range(len(new_prs.slides), 0, -1):
            delete_slide(new_prs, index)

    slide_plans: list[dict[str, Any]] = []
    for i, slide in enumerate(prs_source.slides):
        slide_index = i + 1
        spec, orphans, orphan_names, warnings = _extract_slide(
            slide, slide_index, assets_dir if apply else None
        )
        target_meta, match = _choose_layout(
            template_entry, slide.slide_layout.name, spec
        )
        plan: dict[str, Any] = {
            "slide_index": slide_index,
            "source_layout": slide.slide_layout.name,
            "target_layout": target_meta["name"],
            "layout_match": match,
            "carried_over_as_is": orphan_names,
            "warnings": warnings,
        }
        if apply:
            layout_obj = resolve_layout(new_prs, target_meta["name"])
            mapping = plan_content_mapping(layout_obj, spec)
            apply_add_slide(new_prs, layout_obj, spec, mapping)
            new_slide = new_prs.slides[-1]
            if orphans:
                clone_shape_elements(slide, new_slide, orphans, reassign_ids=True)
            plan["migrated"] = [p["section"] for p in mapping["placements"]]
            plan["unplaced_content"] = mapping["unplaced_content"]
        else:
            # plan against the template's stored schema (no live objects needed)
            roles = [
                ph["role"] for ph in target_meta["placeholders"]
                if ph["role"] not in ("footer", "date", "slide_number")
            ]
            sections = {
                "title": spec.title, "subtitle": spec.subtitle, "body": spec.body,
                "table": spec.table, "image": spec.image,
            }
            wanted_roles = {"title": "title", "subtitle": "subtitle", "body": "body",
                            "table": "table", "image": "picture"}
            migrated, unplaced = [], []
            pool = list(roles)
            for section, value in sections.items():
                if value is None:
                    continue
                role = wanted_roles[section]
                if role in pool:
                    pool.remove(role)
                    migrated.append(section)
                else:
                    unplaced.append(section)
            plan["migrated"] = migrated
            plan["unplaced_content"] = [
                f"{s}: no suitable placeholder in '{target_meta['name']}'"
                for s in unplaced
            ]
        slide_plans.append(plan)

    source_size = (
        emu_to_inches(int(prs_source.slide_width)),
        emu_to_inches(int(prs_source.slide_height)),
    )
    template_size = (
        template_entry["slide_size"]["width_in"],
        template_entry["slide_size"]["height_in"],
    )
    deck_plan = {
        "slide_count": len(slide_plans),
        "slide_size_change": (
            {"from": list(source_size), "to": list(template_size)}
            if source_size != template_size
            else None
        ),
        "slides": slide_plans,
        "risks": _summarize_risks(slide_plans, source_size != template_size),
    }
    return new_prs, deck_plan


def _summarize_risks(slide_plans: list[dict], size_change: bool) -> list[str]:
    risks = []
    scored = [p["slide_index"] for p in slide_plans if p["layout_match"] != "exact"]
    if scored:
        risks.append(
            f"{len(scored)} slide(s) have no exact layout match (slides {scored}); "
            "their layout was chosen by scoring — review them."
        )
    carried = [p["slide_index"] for p in slide_plans if p["carried_over_as_is"]]
    if carried:
        risks.append(
            f"{len(carried)} slide(s) carry freeform/chart shapes over unchanged "
            f"(slides {carried}); their positions may not fit the new design."
        )
    unplaced = [p["slide_index"] for p in slide_plans if p.get("unplaced_content")]
    if unplaced:
        risks.append(f"content could not be placed on slides {unplaced}.")
    if size_change:
        risks.append("slide size changes — carried-over shapes keep absolute positions.")
    return risks

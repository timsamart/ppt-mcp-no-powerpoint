"""Icon library tests: vendored Material set, search, tinted insertion, and
harvesting icons out of PowerPoint files."""

import pytest
from PIL import Image

from ppt_mcp import server
from ppt_mcp.errors import PptMcpError


@pytest.fixture()
def deck(sample_deck):
    deck_id = server.ppt_open_deck(str(sample_deck))["deck_id"]
    yield deck_id
    server.ppt_close_deck(deck_id)


@pytest.fixture()
def icon_source_deck(tmp_path):
    """A deck containing two small icon-like PNGs (one used twice -> dupe),
    and one large photo that must be filtered out."""
    from pptx import Presentation
    from pptx.util import Inches

    icon_a = tmp_path / "gear.png"
    Image.new("RGBA", (64, 64), (200, 30, 30, 255)).save(icon_a)
    icon_b = tmp_path / "arrow.png"
    Image.new("RGBA", (48, 24), (30, 30, 200, 255)).save(icon_b)
    photo = tmp_path / "photo.png"
    Image.new("RGB", (1920, 1080), "gray").save(photo)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(str(icon_a), Inches(1), Inches(1))
    pic.name = "Gear Icon"
    slide.shapes.add_picture(str(icon_a), Inches(2), Inches(1))  # duplicate use
    pic_b = slide.shapes.add_picture(str(icon_b), Inches(3), Inches(1))
    pic_b.name = "Arrow Right"
    slide.shapes.add_picture(str(photo), Inches(1), Inches(3))
    path = tmp_path / "icon_source.pptx"
    prs.save(str(path))
    return path


def test_builtin_material_set_is_vendored():
    sets = server.ppt_list_icon_sets()["icon_sets"]
    material = next(s for s in sets if s["set_id"] == "material")
    assert material["icon_count"] >= 150
    assert material["recolorable"] is True
    assert material["license"] == "Apache-2.0"
    assert material["builtin"] is True


def test_search_finds_risk_icons():
    results = server.ppt_search_icons("risk")["results"]
    assert any(r["icon_id"] == "warning" for r in results)
    results = server.ppt_search_icons("growth")["results"]
    assert results[0]["icon_id"] == "trending_up"


def test_insert_icon_tinted_with_theme_color(deck):
    result = server.ppt_insert_icon(
        deck, 3, "warning", left_in=1.0, top_in=1.0, size_in=0.8, color="accent1"
    )
    assert result["applied"] is True
    assert result["tinted"] is True
    slide = server.ppt_get_slide(deck, 3, response_format="json")
    icon_shape = next(s for s in slide["shapes"] if s["shape_id"] == result["shape_id"])
    assert icon_shape["shape_type"] == "PICTURE"
    assert icon_shape["geometry"]["width_in"] == 0.8
    # rendered PNG is transparent and tinted (not black, not opaque white box)
    theme = server.ppt_extract_theme(deck_id=deck)
    accent = theme["color_scheme"]["accent1"].lstrip("#")
    expected = tuple(int(accent[i : i + 2], 16) for i in (0, 2, 4))
    png = server.icon_library.render_png("material", "warning", f"#{accent}", 96)[0]
    with Image.open(png) as im:
        assert im.mode == "RGBA"
        center = im.getpixel((im.size[0] // 2, int(im.size[1] * 0.8)))
        assert center[:3] == expected  # tinted pixels carry the theme color
        corner = im.getpixel((1, 1))
        assert corner[3] == 0  # transparent background


def test_insert_icon_bad_color_lists_slots(deck):
    with pytest.raises(PptMcpError, match="accent1"):
        server.ppt_insert_icon(deck, 1, "warning", color="brandblue")


def test_harvest_icons_from_pptx(icon_source_deck):
    result = server.ppt_harvest_icons("Corp Icons", path=str(icon_source_deck))
    assert result["set_id"] == "corp_icons"
    assert result["imported_count"] == 2          # photo filtered, dupe collapsed
    assert result["skipped_wrong_size"] == 1
    assert "gear_icon" in result["imported"]
    assert "arrow_right" in result["imported"]

    sets = server.ppt_list_icon_sets()["icon_sets"]
    corp = next(s for s in sets if s["set_id"] == "corp_icons")
    assert corp["builtin"] is False

    found = server.ppt_search_icons("gear", set_id="corp_icons")["results"]
    assert found[0]["icon_id"] == "gear_icon"

    # re-harvesting the same file is a no-op (dedupe by content hash)
    again = server.ppt_harvest_icons("Corp Icons", path=str(icon_source_deck))
    assert again["imported_count"] == 0
    assert again["skipped_duplicates"] >= 2


def test_insert_harvested_icon_as_is(icon_source_deck, deck):
    server.ppt_harvest_icons("Corp Icons", path=str(icon_source_deck))
    result = server.ppt_insert_icon(
        deck, 3, "gear_icon", set_id="corp_icons", left_in=2.0, top_in=2.0, size_in=0.5
    )
    assert result["applied"] is True
    assert result["tinted"] is False  # raster icons are not tinted


def test_harvest_from_open_deck(icon_source_deck):
    deck_id = server.ppt_open_deck(str(icon_source_deck))["deck_id"]
    try:
        result = server.ppt_harvest_icons("From Session", deck_id=deck_id)
        assert result["imported_count"] == 2
    finally:
        server.ppt_close_deck(deck_id)


def test_harvest_requires_exactly_one_source():
    with pytest.raises(PptMcpError, match="exactly one"):
        server.ppt_harvest_icons("X")

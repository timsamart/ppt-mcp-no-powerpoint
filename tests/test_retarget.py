"""M5 tests: apply_template fidelity flow and derived-template extraction."""

import pytest
from pptx.util import Inches

from ppt_mcp import server
from ppt_mcp.models import (
    ChartSeries,
    ChartSpec,
    ContentSpec,
    ImageRef,
    Paragraph,
    Position,
)


@pytest.fixture()
def template_id(sample_deck):
    return server.ppt_register_template(str(sample_deck), name="Target Corp")["template_id"]


@pytest.fixture()
def messy_deck(sample_image, tmp_path):
    """A deck with placeholder content, a picture placeholder, a freeform
    textbox, and a freeform chart — the migration obstacle course."""
    deck = server.ppt_create_deck()["deck_id"]
    server.ppt_add_slide(
        deck, "Title Slide", ContentSpec(title="Annual Report", subtitle="FY26")
    )
    server.ppt_add_slide(
        deck, "Title and Content",
        ContentSpec(
            title="Findings",
            body=[Paragraph(text="Point A"), Paragraph(text="Detail", level=1)],
            notes="remember the details",
        ),
    )
    server.ppt_add_slide(
        deck, "Picture with Caption",
        ContentSpec(title="Vision", image=ImageRef(path=str(sample_image))),
    )
    server.ppt_add_chart(
        deck, 2,
        ChartSpec(chart_type="pie", categories=["A", "B"],
                  series=[ChartSeries(name="s", values=[1, 2])]),
        allow_freeform=True, position=Position(left=5, top=4, width=4, height=3),
    )
    # freeform annotation textbox on slide 1
    prs_path = server.sessions.get(deck).working_path
    from pptx import Presentation

    prs = Presentation(str(prs_path))
    box = prs.slides[0].shapes.add_textbox(Inches(1), Inches(6), Inches(3), Inches(0.5))
    box.text_frame.text = "Confidential"
    prs.save(str(prs_path))
    yield deck
    server.ppt_close_deck(deck)


def test_apply_template_dry_run_reports_fidelity(messy_deck, template_id):
    report = server.ppt_apply_template(messy_deck, template_id)  # dry_run default
    assert report["applied"] is False
    plan = report["plan"]
    assert plan["slide_count"] == 3
    by_index = {p["slide_index"]: p for p in plan["slides"]}
    assert by_index[1]["layout_match"] == "exact"
    assert "Confidential" in str(by_index[1]["carried_over_as_is"])
    assert "title" in by_index[2]["migrated"]
    assert "body" in by_index[2]["migrated"]
    assert by_index[2]["carried_over_as_is"]  # the pie chart
    assert "image" in by_index[3]["migrated"]
    assert plan["risks"]
    # deck untouched
    assert server.ppt_get_deck_overview(messy_deck, response_format="json")["slide_count"] == 3


def test_apply_template_migrates_and_validates(messy_deck, template_id):
    result = server.ppt_apply_template(messy_deck, template_id, dry_run=False)
    assert result["applied"] is True
    assert "validation_summary" in result

    overview = server.ppt_get_deck_overview(messy_deck, response_format="json")
    assert overview["slide_count"] == 3
    assert overview["slides"][0]["title"] == "Annual Report"

    slide2 = server.ppt_get_slide(messy_deck, 2, response_format="json")
    assert slide2["title"] == "Findings"
    body = next(
        s for s in slide2["shapes"]
        if s["placeholder"] and s["placeholder"]["role"] == "body" and s["text"]
    )
    assert [(p["text"], p["level"]) for p in body["text"]] == [
        ("Point A", 0), ("Detail", 1),
    ]
    assert slide2["notes"] == "remember the details"
    assert any("chart" in s for s in slide2["shapes"])  # carried-over pie chart

    slide3 = server.ppt_get_slide(messy_deck, 3, response_format="json")
    assert any(
        s["placeholder"] and s["placeholder"]["role"] == "picture"
        for s in slide3["shapes"]
    )

    # the riskiest tool must be undoable
    server.ppt_undo(messy_deck)
    slide1 = server.ppt_get_slide(messy_deck, 1, response_format="json")
    assert any(
        s["text"] and s["text"][0]["text"] == "Confidential" for s in slide1["shapes"]
    )


def test_apply_template_survives_save_reopen(messy_deck, template_id, tmp_path):
    server.ppt_apply_template(messy_deck, template_id, dry_run=False)
    out = tmp_path / "retargeted.pptx"
    server.ppt_save_deck(messy_deck, str(out))
    reopened = server.ppt_open_deck(str(out))["deck_id"]
    try:
        slide2 = server.ppt_get_slide(reopened, 2, response_format="json")
        assert any("chart" in s for s in slide2["shapes"])  # chart rels remapped
        slide3 = server.ppt_get_slide(reopened, 3, response_format="json")
        assert any(
            s["placeholder"] and s["placeholder"]["role"] == "picture"
            for s in slide3["shapes"]
        )
    finally:
        server.ppt_close_deck(reopened)


def test_extract_template_from_deck(messy_deck):
    result = server.ppt_extract_template_from_deck(messy_deck, name="Recovered Design")
    assert result["derived"] is True
    assert result["layout_usage"]["Title Slide"] == 1
    assert result["layout_usage"]["Title and Content"] == 1
    entry = server.ppt_inspect_template(result["template_id"], response_format="json")
    assert entry["derived"] is True
    assert entry["example_slide_count"] == 3
    # a derived template is immediately usable
    created = server.ppt_create_deck(template_id=result["template_id"])
    try:
        assert created["example_slides_removed"] == 3
    finally:
        server.ppt_close_deck(created["deck_id"])

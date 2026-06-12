"""M3 compliance tests: seed known violations into a clean deck and assert
the validator finds exactly those; conservative repair must reduce them."""

import pytest
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ppt_mcp import server
from ppt_mcp.models import ContentSpec, Paragraph
from ppt_mcp.reader import load_presentation


@pytest.fixture()
def deck(sample_deck):
    deck_id = server.ppt_open_deck(str(sample_deck))["deck_id"]
    yield deck_id
    server.ppt_close_deck(deck_id)


@pytest.fixture()
def template_id(sample_deck):
    return server.ppt_register_template(str(sample_deck), name="Ref")["template_id"]


def _working_prs(deck_id):
    return load_presentation(server.sessions.get(deck_id).working_path)


def _rules(result):
    return {f["rule"] for f in result["findings"]}


def test_clean_deck_has_no_errors(deck, template_id):
    result = server.ppt_validate_compliance(deck, template_id)
    assert result["summary"]["error"] == 0


def test_c02_foreign_font_detected_and_repaired(deck, template_id):
    prs = _working_prs(deck)
    title = prs.slides[0].shapes.title
    title.text_frame.paragraphs[0].runs[0].font.name = "Comic Sans MS"
    prs.save(str(server.sessions.get(deck).working_path))

    result = server.ppt_validate_compliance(deck, template_id)
    c02 = [f for f in result["findings"] if f["rule"] == "C02"]
    assert len(c02) == 1
    assert c02[0]["auto_fixable"] is True
    assert "Comic Sans MS" in c02[0]["message"]

    preview = server.ppt_repair_compliance(deck, template_id)  # dry_run default
    assert preview["applied"] is False
    assert len(preview["planned_fixes"]) == 1

    repaired = server.ppt_repair_compliance(deck, template_id, dry_run=False)
    assert repaired["applied"] is True
    after = server.ppt_validate_compliance(deck, template_id)
    assert "C02" not in _rules(after)


def test_c03_offtheme_color_snap(deck, template_id):
    prs = _working_prs(deck)
    theme_accent = server.ppt_extract_theme(template_id=template_id)["color_scheme"]["accent1"]
    run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    near = RGBColor.from_string(theme_accent.lstrip("#"))
    run.font.color.rgb = RGBColor(min(near[0] + 10, 255), near[1], near[2])
    prs.save(str(server.sessions.get(deck).working_path))

    result = server.ppt_validate_compliance(deck, template_id)
    c03 = [f for f in result["findings"] if f["rule"] == "C03"]
    assert len(c03) == 1 and c03[0]["auto_fixable"] is True

    server.ppt_repair_compliance(deck, template_id, dry_run=False)
    after = server.ppt_validate_compliance(deck, template_id)
    assert "C03" not in _rules(after)
    # color now exactly the theme accent
    prs = _working_prs(deck)
    run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert f"#{run.font.color.rgb}".upper() == theme_accent.upper()


def test_c06_placeholder_bypass(deck, template_id):
    prs = _working_prs(deck)
    slide = prs.slides[1]
    body = slide.placeholders[1]
    body.text_frame.clear()  # empty the body placeholder
    box = slide.shapes.add_textbox(body.left, body.top, body.width, body.height)
    box.text_frame.text = "I bypass the placeholder"
    prs.save(str(server.sessions.get(deck).working_path))

    result = server.ppt_validate_compliance(deck, template_id)
    c06 = [f for f in result["findings"] if f["rule"] == "C06"]
    assert len(c06) == 1
    assert c06[0]["slide_index"] == 2


def test_c07_overflow_estimate(deck, template_id):
    long_bullets = [Paragraph(text=f"Bullet {i}: " + "very long content " * 8) for i in range(14)]
    server.ppt_add_slide(deck, "Title and Content",
                         ContentSpec(title="Overstuffed", body=long_bullets))
    result = server.ppt_validate_compliance(deck, template_id)
    c07 = [f for f in result["findings"] if f["rule"] == "C07"]
    assert any(f["slide_index"] == 4 for f in c07)
    assert all(f["confidence"] == "estimated" for f in c07)
    # C08 density vs layout capacity should fire too
    assert any(f["rule"] == "C08" and f["slide_index"] == 4 for f in result["findings"])


def test_c09_offgrid_shape(deck, template_id):
    prs = _working_prs(deck)
    box = prs.slides[0].shapes.add_textbox(Inches(3.33), Inches(4.77), Inches(2), Inches(0.5))
    box.text_frame.text = "off grid"
    prs.save(str(server.sessions.get(deck).working_path))
    result = server.ppt_validate_compliance(deck, template_id)
    assert "C09" in _rules(result)


def test_validation_without_template_skips_template_rules(deck):
    result = server.ppt_validate_compliance(deck)
    assert result["template_id"] is None
    assert not _rules(result) & {"C01", "C08", "C10"}

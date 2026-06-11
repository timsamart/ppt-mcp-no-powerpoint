"""Shared fixtures. PPT_MCP_HOME is pinned to a temp dir *before* any
ppt_mcp import so the test run never touches the real ~/.ppt-mcp."""

import os
import tempfile
from pathlib import Path

_TEST_HOME = tempfile.mkdtemp(prefix="ppt-mcp-test-home-")
os.environ["PPT_MCP_HOME"] = _TEST_HOME

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_mcp.sessions import SessionManager
from ppt_mcp.store import Store


def build_sample_deck(path: Path) -> Path:
    """Three slides exercising the read path: title slide, bullets with levels
    and notes, and a freeform textbox (non-placeholder shape)."""
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Quarterly Business Review"
    slide1.placeholders[1].text = "FY26 Q2"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Key Risks"
    body = slide2.placeholders[1].text_frame
    body.text = "Unclear data ownership"
    p2 = body.add_paragraph()
    p2.text = "No decision authority"
    p3 = body.add_paragraph()
    p3.text = "Mitigation: governance board"
    p3.level = 1
    slide2.notes_slide.notes_text_frame.text = "Stress the audit trail gap."

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    box = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Freeform annotation"

    prs.save(str(path))
    return path


@pytest.fixture()
def sample_deck(tmp_path: Path) -> Path:
    return build_sample_deck(tmp_path / "sample.pptx")


@pytest.fixture()
def store(tmp_path: Path) -> Store:
    return Store(tmp_path / "ppt-mcp-home")


@pytest.fixture()
def manager(store: Store) -> SessionManager:
    return SessionManager(store)

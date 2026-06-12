"""M2 template intelligence tests: registry, parsing, recommendation,
mapping preview, and template-first deck creation (incl. .potx)."""

import zipfile
from pathlib import Path

import pytest

from ppt_mcp import server
from ppt_mcp.errors import PptMcpError
from ppt_mcp.models import ContentSpec, ImageRef, Paragraph
from ppt_mcp.templates import CT_PRESENTATION, CT_TEMPLATE


@pytest.fixture()
def template_source(sample_deck, tmp_path):
    """A byte-unique copy of the sample deck. Registration is idempotent by
    content hash and the registry is shared across the test session, so
    fixtures that assert on registration metadata must not collide with other
    tests registering identical bytes."""
    import uuid

    from pptx import Presentation

    prs = Presentation(str(sample_deck))
    prs.slides[2].notes_slide.notes_text_frame.text = uuid.uuid4().hex
    path = tmp_path / "unique_template.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def template_id(template_source):
    return server.ppt_register_template(
        str(template_source), name="Test Corporate", version="2.0"
    )["template_id"]


def make_potx(pptx_path: Path, potx_path: Path) -> Path:
    """Turn a .pptx into a real .potx by switching the main content type."""
    with zipfile.ZipFile(pptx_path) as zin, zipfile.ZipFile(
        potx_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(CT_PRESENTATION.encode(), CT_TEMPLATE.encode())
            zout.writestr(item, data)
    return potx_path


def test_register_parses_design_system(template_id):
    entry = server.ppt_inspect_template(template_id, response_format="json")
    assert entry["name"] == "Test Corporate"
    layout_names = [lo["name"] for lo in entry["layouts"]]
    assert "Title Slide" in layout_names
    assert "Comparison" in layout_names
    # theme extraction
    colors = entry["theme"]["color_scheme"]
    assert colors["accent1"].startswith("#")
    assert len(colors) >= 10
    assert entry["theme"]["font_scheme"]["major"]
    # intent tag inference
    by_name = {lo["name"]: lo for lo in entry["layouts"]}
    assert "title_slide" in by_name["Title Slide"]["intent_tags"]
    assert "comparison" in by_name["Comparison"]["intent_tags"]
    assert "image" in by_name["Picture with Caption"]["intent_tags"]
    # capacity estimated for body layouts
    assert by_name["Title and Content"]["capacity"]["body_bullets"] > 3


def test_duplicate_registration_is_idempotent(template_id, template_source):
    again = server.ppt_register_template(str(template_source))
    assert again["template_id"] == template_id
    assert again["already_registered"] is True


def test_inspect_layout_markdown(template_id):
    md = server.ppt_inspect_layout(template_id, "Title and Content")
    assert "role=body" in md
    assert "Capacity" in md


def test_update_template_tags(template_id):
    entry = server.ppt_inspect_template(template_id, response_format="json")
    layout_id = entry["layouts"][1]["layout_id"]
    server.ppt_update_template(
        template_id, {"layout_intent_tags": {layout_id: ["executive_summary"]}}
    )
    updated = server.ppt_inspect_layout(template_id, layout_id, response_format="json")
    assert updated["intent_tags"] == ["executive_summary"]


def test_extract_theme_requires_exactly_one_source(template_id):
    with pytest.raises(PptMcpError, match="exactly one"):
        server.ppt_extract_theme()
    theme = server.ppt_extract_theme(template_id=template_id)
    assert "color_scheme" in theme


def test_recommend_layout_by_name(template_id):
    result = server.ppt_recommend_layout(
        template_id, "comparison of two options",
        ContentSpec(title="A vs B", body=[Paragraph(text="x")]),
    )
    top = result["recommendations"][0]
    assert top["name"] == "Comparison"
    assert top["confidence"] > 0.5
    assert top["reason"]


def test_recommend_layout_structural_image(template_id, sample_image):
    result = server.ppt_recommend_layout(
        template_id, "vision",
        ContentSpec(title="Vision", image=ImageRef(path=str(sample_image))),
    )
    top = result["recommendations"][0]
    assert top["name"] == "Picture with Caption"
    assert top["missing_sections"] == []


def test_map_content_preview_without_deck(template_id):
    plan = server.ppt_map_content_to_placeholders(
        template_id, "Title and Content",
        ContentSpec(title="T", body=[Paragraph(text="b")]),
    )
    assert len(plan["placements"]) == 2
    assert plan["unplaced_content"] == []


def test_create_deck_from_template(template_id):
    created = server.ppt_create_deck(template_id=template_id)
    deck_id = created["deck_id"]
    try:
        assert created["example_slides_removed"] == 3
        overview = server.ppt_get_deck_overview(deck_id, response_format="json")
        assert overview["slide_count"] == 0
        assert "Comparison" in overview["masters"][0]["layouts"]
        result = server.ppt_add_slide(
            deck_id, "Title Slide", ContentSpec(title="From Template")
        )
        assert result["applied"] is True
    finally:
        server.ppt_close_deck(deck_id)


def test_create_deck_keeps_example_slides(template_id):
    created = server.ppt_create_deck(template_id=template_id, include_example_slides=True)
    try:
        overview = server.ppt_get_deck_overview(created["deck_id"], response_format="json")
        assert overview["slide_count"] == 3
    finally:
        server.ppt_close_deck(created["deck_id"])


def test_potx_register_and_create(sample_deck, tmp_path):
    potx = make_potx(sample_deck, tmp_path / "corporate.potx")
    template_id = server.ppt_register_template(str(potx), name="Potx Corp")["template_id"]
    created = server.ppt_create_deck(template_id=template_id)
    deck_id = created["deck_id"]
    try:
        # the materialized working copy must read as a normal presentation
        result = server.ppt_add_slide(
            deck_id, "Title and Content",
            ContentSpec(title="potx works", body=[Paragraph(text="yes")]),
        )
        assert result["applied"] is True
        # and its content type must be presentation, not template
        working = None
        from ppt_mcp.server import sessions as session_mgr
        working = session_mgr.get(deck_id).working_path
        with zipfile.ZipFile(working) as z:
            content_types = z.read("[Content_Types].xml").decode()
        assert CT_PRESENTATION in content_types
        assert CT_TEMPLATE not in content_types
    finally:
        server.ppt_close_deck(deck_id)


def test_unknown_template_error_is_actionable():
    with pytest.raises(PptMcpError, match="ppt_register_template"):
        server.ppt_inspect_template("tpl_nope")
